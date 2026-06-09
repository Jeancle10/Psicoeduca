"""
Genera 3 materiales complementarios:
1. presentacion-casos-clinicos.pptx  — 5 casos para practicar AF (Cap 4)
2. cuestionario-AF.docx              — 20 preguntas sobre AF (Unidad 4)
3. cuestionario-unidad2.docx         — 20 preguntas sobre Bases Filosóficas (Unidad 2)
"""

# ════════════════════════════════════════════════════════
# PARTE 1 — CASOS CLÍNICOS (.pptx)
# ════════════════════════════════════════════════════════
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
from pptx.oxml.ns import qn

DARK_BG = RGBColor(0x1E,0x3A,0x5F); CREMA = RGBColor(0xF2,0xED,0xE4)
NAVY    = RGBColor(0x2B,0x5E,0xA7); SKY   = RGBColor(0x4A,0x9F,0xE0)
ORANGE  = RGBColor(0xE8,0xA8,0x35); GREEN_W = RGBColor(0x4A,0xBF,0xB0)
CREAM_T = RGBColor(0xF5,0xF0,0xDC); DARK_T = RGBColor(0x1E,0x3A,0x5F)
BOX_BG  = RGBColor(0xE4,0xEE,0xF8); GRAY_L = RGBColor(0xCC,0xC8,0xBF)

SERIF='Georgia'; SANS='Calibri'
W=Cm(33.87); H=Cm(19.05); ML=Cm(1.8); CW=Cm(30.27)
LOGO_C = r'C:\Users\MI PC\Desktop\PsicoEduca\Marketing\Rebranding\claro sin fondo.png'
LOGO_D = r'C:\Users\MI PC\Desktop\PsicoEduca\Marketing\Rebranding\oscuro sin fondo.png'
PHOTO  = r'C:\Users\MI PC\Desktop\PsicoEduca\Marketing\BOOK\sentado.JPG'

prs = Presentation(); prs.slide_width=W; prs.slide_height=H
blank = prs.slide_layouts[6]

def bg(sl,c=CREMA):
    f=sl.background.fill; f.solid(); f.fore_color.rgb=c

def hbar(sl,x,y,w,color=ORANGE,t=Pt(3)):
    s=sl.shapes.add_shape(1,x,y,w,t)
    s.fill.solid(); s.fill.fore_color.rgb=color; s.line.fill.background()

def txt(sl,text,x,y,w,h,font=SANS,size=14,bold=False,italic=False,
        color=DARK_T,align=PP_ALIGN.LEFT):
    tb=sl.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text; r.font.name=font; r.font.size=Pt(size)
    r.font.bold=bold; r.font.italic=italic; r.font.color.rgb=color
    return tb

def paras(sl,lines,x,y,w,h,font=SANS,size=13,color=DARK_T,align=PP_ALIGN.LEFT):
    tb=sl.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    for i,ln in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if isinstance(ln,dict):
            r=p.add_run(); r.text=ln.get('t','')
            r.font.name=ln.get('font',font); r.font.size=Pt(ln.get('size',size))
            r.font.bold=ln.get('bold',False); r.font.italic=ln.get('italic',False)
            r.font.color.rgb=ln.get('color',color)
        else:
            r=p.add_run(); r.text=str(ln); r.font.name=font
            r.font.size=Pt(size); r.font.color.rgb=color

def logo(sl,dark=False):
    path=LOGO_D if dark else LOGO_C
    try: sl.shapes.add_picture(path,W-Cm(4.5),Cm(0.3),Cm(4.0),Cm(2.2))
    except: pass

def footer(sl,dark=False):
    c=CREAM_T if dark else RGBColor(0x88,0x88,0x88)
    txt(sl,'Lic. Jean Clemotte  |  @Psico_Educa20',ML,H-Cm(1.1),Cm(22),Cm(0.8),
        SANS,10,italic=True,color=c)

def hmod(sl,label):
    txt(sl,label.upper(),ML,Cm(0.5),Cm(22),Cm(0.7),SANS,9,bold=True,color=NAVY)
    hbar(sl,ML,Cm(1.3),Cm(22),NAVY,Pt(1))

def box(sl,text,x,y,w,h,fill=BOX_BG,border=NAVY,font=SANS,size=13,
        bold=False,color=DARK_T,align=PP_ALIGN.LEFT):
    s=sl.shapes.add_shape(5,x,y,w,h)
    s.fill.solid(); s.fill.fore_color.rgb=fill
    s.line.color.rgb=border; s.line.width=Pt(1)
    tf=s.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text; r.font.name=font; r.font.size=Pt(size)
    r.font.bold=bold; r.font.color.rgb=color; return s

# ── S1 PORTADA ───────────────────────────────────────────
sl=prs.slides.add_slide(blank); bg(sl,DARK_BG)
txt(sl,'Casos Clínicos',ML,Cm(2.0),Cm(22),Cm(3.5),SERIF,54,bold=True,color=CREAM_T)
txt(sl,'Práctica del Análisis Funcional paso a paso',ML,Cm(6.0),Cm(25),Cm(1.5),SANS,18,color=SKY)
hbar(sl,ML,Cm(7.8),Cm(20),ORANGE,Pt(3))
txt(sl,'5 casos clínicos contextualizados — Paraguay',ML,Cm(8.5),Cm(22),Cm(1.2),
    SANS,15,italic=True,color=RGBColor(0xAA,0xBB,0xCC))
txt(sl,'Unidad 4 — Froxán et al.',ML,Cm(10.2),Cm(18),Cm(1.1),
    SERIF,13,italic=True,color=RGBColor(0x88,0x99,0xAA))
txt(sl,'↓↓↓',W-Cm(9),Cm(2.5),Cm(7),Cm(13),SANS,130,bold=True,
    color=RGBColor(0x2B,0x4A,0x6E),align=PP_ALIGN.CENTER)
logo(sl,dark=True); footer(sl,dark=True)

# ── S2 INSTRUCCIONES + GUÍA AF ───────────────────────────
sl=prs.slides.add_slide(blank); bg(sl,CREMA)
hmod(sl,'Casos Clínicos — Guía de análisis')
txt(sl,'Procedimiento: Análisis Funcional paso a paso',ML,Cm(1.7),CW-Cm(5),Cm(1.8),
    SERIF,28,bold=True,color=DARK_T)

pasos=[
    ('PASO 1','Análisis morfológico — ¿Qué hace exactamente?',
     'Describí la conducta problema en términos operativos: frecuencia, duración y localización temporal.\n❌ No usar etiquetas ("tiene ansiedad"). ✓ Describir qué hace, cuándo y con qué frecuencia.'),
    ('PASO 2','Identificar antecedentes y consecuentes',
     'Antecedentes: ¿qué ocurre ANTES? (¿dónde está, con quién, qué piensa, qué siente?)\nConsecuentes: ¿qué ocurre DESPUÉS? (¿qué obtiene, qué evita, cómo reacciona el entorno?)'),
    ('PASO 3','Variables disposicionales y motivadoras',
     'Disposicionales (estables): edad, historia de aprendizaje, reglas, entorno social.\nMotivadoras (cambian): estado emocional, privación, anticipaciones verbales.'),
    ('PASO 4','Hipótesis funcional',
     'Origen: ¿cómo se aprendió? / Mantenimiento: ¿qué la sostiene HOY?\nIdentificá el procedimiento operante que mantiene la conducta (RF+, RN, CP, CN).'),
]
bw=CW/2-Cm(0.6)
for i,(num,tit,desc) in enumerate(pasos):
    col=i%2; row=i//2
    x=ML+col*(bw+Cm(0.6)); y=Cm(4.2)+row*Cm(5.8)
    s=sl.shapes.add_shape(1,x,y,bw,Cm(5.3))
    s.fill.solid(); s.fill.fore_color.rgb=BOX_BG
    accent=ORANGE if col==0 else GREEN_W
    s.line.color.rgb=accent; s.line.width=Pt(2)
    txt(sl,num,x+Cm(0.3),y+Cm(0.2),Cm(2),Cm(0.9),SERIF,14,bold=True,color=accent)
    txt(sl,tit,x+Cm(0.3),y+Cm(1.1),bw-Cm(0.6),Cm(0.9),SANS,13,bold=True,color=DARK_T)
    txt(sl,desc,x+Cm(0.3),y+Cm(2.0),bw-Cm(0.6),Cm(3.0),SANS,11,color=DARK_T)

logo(sl); footer(sl)

# ── CASOS ────────────────────────────────────────────────
casos=[
    {
        'num':'01','personaje':'Marge Simpson','edad':'39 años, ama de casa, Asunción',
        'motivo':'"Homero la trajo porque no puede salir sola a ningún lado."',
        'datos':[
            '• Hace 2 años fue asaltada en el Super Stock de Villa Morra.',
            '• Desde entonces evita ir sola a supermercados, shoppings y reuniones.',
            '• Cuando Homero le propone salir, dice "me duele la cabeza" o "estoy cansada".',
            '• Si sale sola, a los 10 minutos: taquicardia, sensación de ahogo, regresa.',
            '• Con Homero presente, puede ir a cualquier lugar sin dificultad.',
            '• Homero fue asumiendo todas las compras y salidas sin quejarse.',
            '• Regla expresada: "Si salgo sola, me van a asaltar de nuevo."',
            '• El estrés del hogar (Bart, Maggie, trabajo de Homero) agudiza los síntomas.',
        ],
        'preguntas':[
            '1. ¿Cuál es la conducta problema? Describila operativamente (frecuencia, duración, localización).',
            '2. Identificá los antecedentes (Ed) y consecuentes (C) de la conducta de evitación.',
            '3. ¿Qué respuesta condicionada (E-R) está presente? Identificá EI, EN/EC, RI, RC.',
            '4. Listá las variables disposicionales y motivadoras. Especificá cómo afectan la contingencia.',
            '5. Formulá la hipótesis de origen y la hipótesis de mantenimiento.',
        ],
        'color_dark': RGBColor(0x4A,0x2A,0x6A),
    },
    {
        'num':'02','personaje':'Bart Simpson','edad':'10 años, estudiante, Asunción',
        'motivo':'"Sus padres lo traen porque no puede concentrarse y hace berrinches para no estudiar."',
        'datos':[
            '• En clase de matemáticas: interrumpe, hace bromas, tira útiles si la maestra lo llama al frente.',
            '• En clase de música (no tiene notas): trabaja bien, levanta la mano, está relajado.',
            '• Antes de los exámenes: dice que le duele la panza y pide quedarse en casa.',
            '• En casa, cuando le piden que estudie: llora hasta que sus padres "le dan un descanso".',
            '• Al ceder, los padres le dan el teléfono "para que se calme".',
            '• Historia: el año pasado la maestra lo ridiculizó delante de toda la clase.',
            '• Cuando falta a la clase de mate, está tranquilo el resto del día.',
        ],
        'preguntas':[
            '1. Describí operativamente las conductas problema (hay más de una).',
            '2. Identificá qué procedimientos operantes mantienen las conductas de Bart (RF+, RN, CP, CN).',
            '3. ¿Cuál es el Ed+ y el Ed- para la conducta disruptiva en clase?',
            '4. ¿Qué papel juega el episodio de ridiculización en la hipótesis de ORIGEN?',
            '5. Formulá la hipótesis de mantenimiento y proponé el objetivo del tratamiento.',
        ],
        'color_dark': RGBColor(0x1E,0x4A,0x2A),
    },
    {
        'num':'03','personaje':'Homero Simpson','edad':'39 años, empleado Planta Nuclear, Asunción',
        'motivo':'"Marge lo trae porque come demasiado y no colabora en la casa."',
        'datos':[
            '• Come donuts y papas fritas principalmente al llegar del trabajo y cuando hay tareas en casa.',
            '• Si hay tareas domésticas pendientes, va directo a la heladera antes de que Marge pida algo.',
            '• En el trabajo, durante la jornada laboral, come mucho menos.',
            '• Cuando come, dice que siente que "todo es más llevadero".',
            '• Marge dejó de pedirle ayuda porque "siempre termina en discusión".',
            '• Historia: su madre le daba comida como premio o "para animarlo" cuando estaba angustiado.',
            '• Regla expresada: "Después de trabajar todo el día, me lo merezco."',
        ],
        'preguntas':[
            '1. Describí la conducta problema: ¿qué hace Homero exactamente? (frecuencia, localización).',
            '2. ¿Qué tipo de reforzamiento mantiene la conducta de comer? ¿Y la de no colaborar?',
            '3. ¿Qué función cumple la regla "me lo merezco"? ¿Es Ed, variable disposicional o motivadora?',
            '4. ¿Qué papel juega la historia de aprendizaje (madre-comida) como variable disposicional?',
            '5. Formulá hipótesis de mantenimiento e identificá qué mantiene el problema de Marge también.',
        ],
        'color_dark': RGBColor(0x4A,0x3A,0x1A),
    },
    {
        'num':'04','personaje':'Lisa Simpson','edad':'8 años, estudiante, Asunción',
        'motivo':'"Llora mucho antes de los exámenes y dice que no puede dormir."',
        'datos':[
            '• La noche antes de cualquier examen: no duerme, llora, dice "no me sé nada".',
            '• Por las mañanas de examen: náuseas o dolor de estómago.',
            '• En el colegio: levanta la mano en cada pregunta; si no la llaman, se pone ansiosa.',
            '• Si saca un 9 en lugar de 10, llora cuando llega a casa.',
            '• Sus padres le dicen constantemente: "Vos podés ser la mejor, sos la más inteligente."',
            '• Cuando llora, su madre la abraza y a veces la excusa del colegio si "está muy mal".',
            '• En actividades sin calificaciones (arte, música): completamente relajada y disfruta.',
        ],
        'preguntas':[
            '1. Identificá al menos 3 conductas problema diferentes. Describilas operativamente.',
            '2. ¿Qué respuesta condicionada está presente? Identificá los elementos del CC.',
            '3. ¿Qué consecuentes mantienen las conductas de ansiedad? (RF+, RN, CP, CN)',
            '4. ¿Qué variables disposicionales están operando? ¿Cuáles son motivadoras?',
            '5. ¿Cómo el comportamiento de los padres actúa como consecuente? ¿Qué implicaciones tiene para el tratamiento?',
        ],
        'color_dark': RGBColor(0x1A,0x3A,0x4A),
    },
    {
        'num':'05','personaje':'Carlos (caso paraguayo)','edad':'28 años, contador, Asunción',
        'motivo':'"Tengo ansiedad. No puedo dar presentaciones en el trabajo."',
        'datos':[
            '• Días previos a presentar ante jefes: no duerme bien, come poco, pensamientos repetitivos.',
            '• Día de la presentación: taquicardia, manos sudadas, voz temblorosa.',
            '• Ha pedido "cambiar" sus presentaciones con colegas varias veces ("estoy ocupado").',
            '• Cuando logra evitar una presentación, siente alivio inmediato y duerme bien esa noche.',
            '• Cuando inevitablemente presenta, termina sin problema y sus colegas lo felicitan.',
            '• En contextos informales (asado, fútbol, chopp con amigos): "es el más payaso del grupo".',
            '• Historia: en la universidad, un profesor lo ridiculizó públicamente durante una exposición.',
        ],
        'preguntas':[
            '1. Describí la conducta problema operativamente. ¿Es una o son varias? ¿Hay conducta E-R y E-R-C?',
            '2. Identificá el EI, EC, RI y RC. ¿Cuándo ocurrió el condicionamiento?',
            '3. ¿Qué mantiene la conducta de evitación? ¿Y qué la debilitaría?',
            '4. ¿Qué variables disposicionales y motivadoras están presentes en el caso de Carlos?',
            '5. Formulá la hipótesis de mantenimiento y proponé el foco del tratamiento.',
        ],
        'color_dark': RGBColor(0x3A,0x1A,0x4A),
    },
]

for caso in casos:
    # ── SLIDE A: DATOS DEL CASO (DARK con acento de color)
    sl=prs.slides.add_slide(blank); bg(sl,DARK_BG)

    # Header badge del caso
    s=sl.shapes.add_shape(5,ML,Cm(1.5),Cm(3.5),Cm(1.8))
    s.fill.solid(); s.fill.fore_color.rgb=ORANGE; s.line.fill.background()
    tf=s.text_frame; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text='CASO '+caso['num']
    r.font.name=SERIF; r.font.size=Pt(16); r.font.bold=True; r.font.color.rgb=DARK_BG

    txt(sl,caso['personaje'],ML+Cm(4.2),Cm(1.5),CW-Cm(9),Cm(1.0),
        SERIF,26,bold=True,color=CREAM_T)
    txt(sl,caso['edad'],ML+Cm(4.2),Cm(2.7),CW-Cm(9),Cm(0.8),
        SANS,14,color=SKY)

    hbar(sl,ML,Cm(3.7),CW-Cm(5),ORANGE,Pt(2))

    txt(sl,'Motivo de consulta:',ML,Cm(4.2),CW-Cm(5),Cm(0.8),
        SANS,12,bold=True,color=ORANGE)
    txt(sl,caso['motivo'],ML,Cm(5.1),CW-Cm(5),Cm(1.1),
        SANS,14,italic=True,color=CREAM_T)

    txt(sl,'Información relevante:',ML,Cm(6.5),CW-Cm(5),Cm(0.8),
        SANS,12,bold=True,color=GREEN_W)
    paras(sl,caso['datos'],ML,Cm(7.3),CW-Cm(5),Cm(9.0),SANS,13,CREAM_T)

    logo(sl,dark=True); footer(sl,dark=True)

    # ── SLIDE B: PREGUNTAS DE ANÁLISIS (CREMA)
    sl=prs.slides.add_slide(blank); bg(sl,CREMA)
    hmod(sl,'Casos Clínicos — Análisis Funcional paso a paso')

    # Header
    s2=sl.shapes.add_shape(5,ML,Cm(1.7),Cm(3.5),Cm(1.6))
    s2.fill.solid(); s2.fill.fore_color.rgb=DARK_BG; s2.line.fill.background()
    tf2=s2.text_frame; p2=tf2.paragraphs[0]; p2.alignment=PP_ALIGN.CENTER
    r2=p2.add_run(); r2.text='CASO '+caso['num']
    r2.font.name=SERIF; r2.font.size=Pt(15); r2.font.bold=True; r2.font.color.rgb=ORANGE

    txt(sl,caso['personaje']+' — Guía de análisis',
        ML+Cm(4.0),Cm(1.7),CW-Cm(8),Cm(0.9),
        SERIF,20,bold=True,color=DARK_T)
    txt(sl,'Aplicá el procedimiento AF paso a paso:',
        ML+Cm(4.0),Cm(2.7),CW-Cm(8),Cm(0.8),
        SANS,13,italic=True,color=NAVY)

    hbar(sl,ML,Cm(3.6),CW-Cm(1),ORANGE,Pt(2))

    for j,preg in enumerate(caso['preguntas']):
        y=Cm(4.2)+j*Cm(2.6)
        accent=ORANGE if j%2==0 else GREEN_W
        s3=sl.shapes.add_shape(5,ML,y,Cm(0.7),Cm(2.3))
        s3.fill.solid(); s3.fill.fore_color.rgb=accent; s3.line.fill.background()
        tf3=s3.text_frame; p3=tf3.paragraphs[0]; p3.alignment=PP_ALIGN.CENTER
        r3=p3.add_run(); r3.text=str(j+1)
        r3.font.name=SERIF; r3.font.size=Pt(14); r3.font.bold=True; r3.font.color.rgb=CREAM_T
        box(sl,preg,ML+Cm(0.9),y,CW-Cm(1.8),Cm(2.3),
            fill=BOX_BG,border=RGBColor(0xCC,0xC8,0xC0),size=13,color=DARK_T)

    logo(sl); footer(sl)

# ── CIERRE ───────────────────────────────────────────────
sl=prs.slides.add_slide(blank)
try: sl.shapes.add_picture(PHOTO,0,0,W,H)
except: pass
from pptx.oxml.ns import qn as _q; from lxml import etree as _e
ov=sl.shapes.add_shape(1,0,0,W*0.54,H)
ov.fill.solid(); ov.fill.fore_color.rgb=DARK_BG; ov.line.fill.background()
xPr=ov.fill._xPr; sf=xPr.solidFill
c=sf.find(_q('a:srgbClr'))
if c is None: c=_e.SubElement(sf,_q('a:srgbClr')); c.set('val','1E3A5F')
a=_e.SubElement(c,_q('a:alpha')); a.set('val','85000')
hbar(sl,ML,Cm(1.8),Cm(16),ORANGE,Pt(2))
hbar(sl,ML,Cm(17.0),Cm(16),ORANGE,Pt(2))
txt(sl,'Gracias por\nsu Atención',ML,Cm(3.5),Cm(16),Cm(7),
    SERIF,50,bold=True,italic=True,color=CREAM_T)
txt(sl,'Casos Clínicos — Análisis Funcional',ML,Cm(11.5),Cm(16),Cm(2),SANS,16,color=SKY)
txt(sl,'Lic. Jean Clemotte  |  PsicoEduca',ML,Cm(13.8),Cm(16),Cm(1.3),
    SANS,14,italic=True,color=CREAM_T)
logo(sl,dark=True)

out1=r'C:\Users\MI PC\psicoeduca\materiales\presentacion-casos-clinicos.pptx'
prs.save(out1); print('1/3 OK: '+out1)


# ════════════════════════════════════════════════════════
# PARTE 2 — CUESTIONARIOS (.docx)
# ════════════════════════════════════════════════════════
from docx import Document
from docx.shared import Pt as DPt, RGBColor as DRGB, Cm as DCm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn as dqn
from docx.oxml import OxmlElement

def make_quiz(filename, titulo, instrucciones, preguntas, respuestas):
    doc = Document()

    # Márgenes
    for section in doc.sections:
        section.top_margin    = DCm(2.0)
        section.bottom_margin = DCm(2.0)
        section.left_margin   = DCm(2.5)
        section.right_margin  = DCm(2.5)

    # ── ENCABEZADO ────────────────────────────────────────
    h = doc.add_heading(titulo, 0)
    h.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in h.runs:
        run.font.color.rgb = DRGB(0x1E,0x3A,0x5F)
        run.font.size = DPt(20)

    p_sub = doc.add_paragraph()
    p_sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p_sub.add_run('Lic. Jean Clemotte  |  @Psico_Educa20  |  PsicoEduca')
    run.font.size = DPt(10); run.italic = True
    run.font.color.rgb = DRGB(0x66,0x66,0x66)

    doc.add_paragraph()

    # ── INSTRUCCIONES ─────────────────────────────────────
    p_inst = doc.add_paragraph()
    r1 = p_inst.add_run('Instrucciones: ')
    r1.bold = True; r1.font.color.rgb = DRGB(0x1E,0x3A,0x5F); r1.font.size = DPt(11)
    r2 = p_inst.add_run(instrucciones)
    r2.font.size = DPt(11)

    doc.add_paragraph()

    # ── DATOS DEL ESTUDIANTE ──────────────────────────────
    p_datos = doc.add_paragraph()
    p_datos.add_run('Nombre: ______________________________________   '
                    'Fecha: _______________   Nota: _______')
    doc.add_paragraph()

    # ── PREGUNTAS ─────────────────────────────────────────
    opciones = ['A)', 'B)', 'C)', 'D)']
    for i, (pregunta, opts) in enumerate(preguntas):
        # Número + pregunta
        p_q = doc.add_paragraph()
        r_num = p_q.add_run(f'{i+1}. ')
        r_num.bold = True; r_num.font.color.rgb = DRGB(0x1E,0x3A,0x5F)
        r_num.font.size = DPt(11)
        r_txt = p_q.add_run(pregunta)
        r_txt.font.size = DPt(11); r_txt.bold = True

        # Opciones
        for j, opt in enumerate(opts):
            p_opt = doc.add_paragraph(style='List Bullet')
            p_opt.paragraph_format.left_indent = DCm(1.0)
            r_opt = p_opt.add_run(f'{opciones[j]}  {opt}')
            r_opt.font.size = DPt(11)

        doc.add_paragraph()

    # ── SALTO DE PÁGINA antes de la matriz ────────────────
    doc.add_page_break()

    # ── MATRIZ DE RESPUESTAS ──────────────────────────────
    h_matriz = doc.add_heading('Matriz de Respuestas Correctas', 1)
    h_matriz.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in h_matriz.runs:
        run.font.color.rgb = DRGB(0x1E,0x3A,0x5F)

    doc.add_paragraph()

    p_nota = doc.add_paragraph()
    r_nota = p_nota.add_run(
        'Nota: Esta página es para uso exclusivo del docente. '
        'No distribuir junto con el cuestionario del estudiante.')
    r_nota.italic = True; r_nota.font.size = DPt(10)
    r_nota.font.color.rgb = DRGB(0x88,0x88,0x88)
    doc.add_paragraph()

    # Tabla de respuestas 4 columnas × 5 filas = 20 respuestas
    table = doc.add_table(rows=6, cols=4)
    table.style = 'Table Grid'

    headers = ['Preg.', 'Respuesta', 'Preg.', 'Respuesta']
    hrow = table.rows[0]
    for ci, hdr in enumerate(headers):
        cell = hrow.cells[ci]
        cell.text = hdr
        for para in cell.paragraphs:
            for run in para.runs:
                run.bold = True; run.font.size = DPt(11)
                run.font.color.rgb = DRGB(0xFF,0xFF,0xFF)
        # Fondo navy para encabezado
        tc = cell._tc; tcPr = tc.get_or_add_tcPr()
        shd = OxmlElement('w:shd')
        shd.set(dqn('w:val'), 'clear')
        shd.set(dqn('w:color'), 'auto')
        shd.set(dqn('w:fill'), '1E3A5F')
        tcPr.append(shd)

    for ri in range(5):
        row = table.rows[ri+1]
        for ci in range(4):
            if ci == 0:
                q_num = ri + 1
                row.cells[ci].text = str(q_num)
                row.cells[ci+1].text = respuestas[q_num-1]
            elif ci == 2:
                q_num = ri + 6
                row.cells[ci].text = str(q_num)
                row.cells[ci+1].text = respuestas[q_num-1]
        for cell in row.cells:
            for para in cell.paragraphs:
                para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                for run in para.runs:
                    run.font.size = DPt(11)

    # Segunda tabla filas 11-20
    doc.add_paragraph()
    table2 = doc.add_table(rows=6, cols=4)
    table2.style = 'Table Grid'
    hrow2 = table2.rows[0]
    for ci, hdr in enumerate(headers):
        cell = hrow2.cells[ci]; cell.text = hdr
        for para in cell.paragraphs:
            for run in para.runs:
                run.bold = True; run.font.size = DPt(11)
                run.font.color.rgb = DRGB(0xFF,0xFF,0xFF)
        tc = cell._tc; tcPr = tc.get_or_add_tcPr()
        shd = OxmlElement('w:shd')
        shd.set(dqn('w:val'),'clear'); shd.set(dqn('w:color'),'auto')
        shd.set(dqn('w:fill'),'1E3A5F'); tcPr.append(shd)

    for ri in range(5):
        row2 = table2.rows[ri+1]
        for ci in range(4):
            if ci == 0:
                q_num = ri + 11
                row2.cells[ci].text = str(q_num)
                row2.cells[ci+1].text = respuestas[q_num-1]
            elif ci == 2:
                q_num = ri + 16
                row2.cells[ci].text = str(q_num)
                row2.cells[ci+1].text = respuestas[q_num-1]
        for cell in row2.cells:
            for para in cell.paragraphs:
                para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                for run in para.runs:
                    run.font.size = DPt(11)

    doc.save(filename)


# ── CUESTIONARIO AF (UNIDAD 4) ────────────────────────────
preguntas_AF = [
    ("¿Cuál es el objetivo fundamental del Análisis Funcional de la Conducta?",
     ["Diagnosticar el trastorno psicológico del paciente",
      "Establecer relaciones de contingencia entre estímulos y respuestas para explicar la conducta",
      "Identificar los estados mentales que causan la conducta",
      "Clasificar la conducta según su morfología observable"]),

    ("Según el AF, ¿qué define a una 'clase de respuesta'?",
     ["Respuestas con la misma morfología (aspecto físico)",
      "Respuestas con la misma función en su contexto, aunque difieran en su forma",
      "Las respuestas que ocurren con mayor frecuencia en el repertorio del individuo",
      "Las respuestas que el terapeuta considera problemáticas"]),

    ("En el condicionamiento clásico, ¿cuál de los siguientes es correctamente un EI (Estímulo Incondicionado)?",
     ["Un semáforo en rojo que evoca ansiedad en alguien que tuvo un accidente",
      "Un asalto que provoca miedo de manera innata en el organismo",
      "El consultorio del dentista que genera ansiedad anticipatoria",
      "El sonido de la ambulancia que genera alarma"]),

    ("¿Cómo se define la habituación?",
     ["El aumento de la respuesta ante la repetición del estímulo",
      "El aprendizaje de una nueva conducta por condicionamiento operante",
      "El decremento o desaparición de una respuesta por la repetición del estímulo que la desencadena",
      "La generalización de una respuesta condicionada a nuevos estímulos similares"]),

    ("Carlos fue mordido por un perro labrador. Hoy siente ansiedad ante cualquier perro, incluso de otras razas. ¿Qué fenómeno describe esto?",
     ["Sensibilización al estímulo condicionado",
      "Extinción pavloviana de la respuesta condicionada",
      "Generalización del estímulo en el condicionamiento clásico",
      "Reforzamiento negativo de la conducta de evitación"]),

    ("¿Cuál es la definición correcta de reforzamiento negativo?",
     ["Un estímulo desagradable presentado contingentemente que disminuye la probabilidad de la respuesta",
      "La retirada de un reforzador positivo contingente que disminuye la probabilidad de la respuesta",
      "Un procedimiento que incrementa la probabilidad de la respuesta al retirar un estímulo aversivo",
      "La ausencia de consecuencias contingentes a una respuesta (acontingencia)"]),

    ("Marge evita ir al Super Stock porque siente taquicardia. Al evitar, la taquicardia desaparece. ¿Qué procedimiento mantiene su conducta de evitación?",
     ["Reforzamiento positivo (aparece un reforzador positivo)",
      "Reforzamiento negativo (desaparece el estímulo aversivo)",
      "Castigo positivo (aparece un estímulo aversivo)",
      "Castigo negativo (desaparece un reforzador positivo)"]),

    ("¿Cuál de las siguientes afirmaciones describe correctamente un Estímulo Discriminativo (Ed)?",
     ["Es cualquier estímulo que ocurre temporalmente antes de la respuesta",
      "Es un concepto descriptivo que indica solo relación temporal con la conducta",
      "Es un concepto funcional que señala la disponibilidad diferencial de un reforzador",
      "Es un estímulo que elicita una respuesta de manera innata e incondicionada"]),

    ("Un estudiante levanta la mano en clase. El profesor lo felicita. Desde entonces levanta la mano con más frecuencia. ¿Qué procedimiento describe esto?",
     ["Castigo positivo: aparece un estímulo aversivo tras la respuesta",
      "Reforzamiento negativo: desaparece un estímulo aversivo tras la respuesta",
      "Reforzamiento positivo: aparece un reforzador contingente que aumenta la R",
      "Extinción operante: no hay consecuencias contingentes"]),

    ("¿Cuál es la diferencia entre el antecedente y el estímulo discriminativo?",
     ["El antecedente es funcional; el discriminativo es descriptivo",
      "El antecedente ocurre después de la conducta; el discriminativo ocurre antes",
      "El antecedente es un concepto descriptivo (relación temporal); el Ed es funcional (señala contingencias)",
      "No hay diferencia relevante; son sinónimos en el análisis de conducta"]),

    ("¿Qué son las variables disposicionales?",
     ["Variables que cambian momento a momento y afectan temporalmente el valor del reforzador",
      "Condiciones estables del organismo y el entorno que modifican el valor de los elementos de la contingencia",
      "Los reforzadores positivos disponibles en el contexto del individuo",
      "Las respuestas encubiertas que acompañan a una conducta manifiesta"]),

    ("Homero lleva 3 días sin dormir por el trabajo. En ese contexto, cualquier crítica de Marge lo enoja más de lo habitual. El cansancio actúa como:",
     ["Estímulo discriminativo positivo (Ed+)",
      "Reforzador condicionado secundario",
      "Variable motivadora (operación de establecimiento transitoria)",
      "Variable disposicional biológica permanente"]),

    ("¿Cuál es el primer paso para realizar un análisis funcional en contextos naturales?",
     ["Identificar las variables disposicionales del caso clínico",
      "Formular la hipótesis de mantenimiento de la conducta",
      "Realizar el análisis morfológico: describir operativamente la conducta problema",
      "Identificar el estímulo discriminativo que controla la respuesta"]),

    ("Al describir morfológicamente una conducta, se deben evaluar tres dimensiones. ¿Cuáles son?",
     ["Pensamientos asociados, historia familiar y nivel socioeconómico",
      "Frecuencia de aparición, extensión temporal y localización temporal",
      "Diagnóstico DSM, adherencia al tratamiento y motivación del paciente",
      "Solo la topografía de la conducta observable por el terapeuta"]),

    ("¿Cuál es la diferencia entre la hipótesis de origen y la hipótesis de mantenimiento?",
     ["La hipótesis de origen es la más importante para diseñar el tratamiento",
      "La hipótesis de mantenimiento se refiere al pasado y la de origen al presente",
      "La hipótesis de origen explica cómo se aprendió la conducta; la de mantenimiento explica qué la sostiene HOY",
      "Son sinónimas; ambas explican por qué la conducta se mantiene actualmente"]),

    ("Según el AF, ¿qué caracteriza a un comportamiento como 'problemático'?",
     ["Que su morfología sea socialmente inaceptable o anormal",
      "Que el paciente la reconozca explícitamente como un problema",
      "Que se desvíe de lo esperable en términos de adaptación al entorno u objetivos personales a largo plazo",
      "Que el manual DSM la clasifique como criterio de un trastorno mental"]),

    ("Un cliente llega y dice: 'Tengo depresión'. ¿Cuál es el primer paso correcto desde el AF?",
     ["Aceptar el diagnóstico y diseñar el tratamiento para depresión",
      "Derivar a psiquiatría para confirmar y complementar el diagnóstico",
      "Obtener una descripción operativa de las conductas concretas que incluye el término 'depresión' para ese cliente",
      "Explorar la historia familiar y traumas del pasado para establecer el origen"]),

    ("La evitación de Marge de lugares concurridos es mantenida por reforzamiento negativo. ¿Qué implica esto para el tratamiento?",
     ["Hay que castigar directamente la conducta de evitación",
      "Hay que encontrar y modificar las contingencias que mantienen la evitación (y su alivio)",
      "Hay que aumentar los reforzadores positivos en casa para compensar",
      "Hay que trabajar exclusivamente con la regla verbal 'si salgo, me van a asaltar'"]),

    ("¿Cuál de los siguientes ilustra correctamente la distinción entre morfología y funcionalidad?",
     ["Dos conductas con idéntica morfología siempre tienen la misma función operante",
      "Pedir dinero, llorar y amenazar pueden pertenecer a la misma clase operante si tienen la misma función",
      "La función de una conducta se determina por su topografía y características físicas",
      "La morfología es más importante que la función para realizar un análisis conductual completo"]),

    ("En contextos naturales, ¿por qué el analista de conducta no puede hablar de 'análisis funcional propiamente dicho'?",
     ["Porque no dispone de la historia clínica completa del paciente",
      "Porque no es posible la manipulación experimental necesaria para verificar la función hipotetizada",
      "Porque los contextos naturales no permiten identificar estímulos antecedentes relevantes",
      "Porque el análisis funcional solo puede realizarse en laboratorio con animales"]),
]

respuestas_AF = ['B','B','B','C','C','C','B','C','C','C','B','C','C','B','C','C','C','B','B','B']

make_quiz(
    r'C:\Users\MI PC\psicoeduca\materiales\cuestionario-AF.docx',
    'Cuestionario — Análisis Funcional de la Conducta (Unidad 4)',
    'Seleccioná la única opción correcta para cada pregunta. '
    'Cada pregunta vale 0,5 puntos. Total: 10 puntos.',
    preguntas_AF,
    respuestas_AF
)
print('2/3 OK: cuestionario-AF.docx')


# ── CUESTIONARIO UNIDAD 2 (BASES FILOSÓFICAS) ─────────────
preguntas_U2 = [
    ("¿Qué es el conductismo radical en el contexto del análisis de la conducta?",
     ["Una teoría que niega totalmente la existencia de los estados mentales",
      "La filosofía de la psicología skinneriana que considera la conducta como raíz de lo psicológico",
      "Un método terapéutico basado exclusivamente en el condicionamiento operante",
      "Una corriente que solo estudia conductas observables por terceros"]),

    ("¿Qué denomina Sellars (1956) 'imagen científica'?",
     ["La explicación del comportamiento en términos de creencias, deseos e intenciones",
      "El diagnóstico clínico basado en el DSM",
      "La explicación del comportamiento en términos de causas materiales (estados cerebrales, contingencias)",
      "El nivel de análisis normativo propio de la ética y la filosofía moral"]),

    ("Según Ryle (2005), ¿en qué consiste el 'error categorial' en la filosofía de Descartes?",
     ["En separar mente y cuerpo como dos sustancias de naturaleza completamente distinta",
      "En concebir lo mental como una entidad fáctica (una 'cosa') que mantiene relaciones causales con el cuerpo",
      "En negar que los estados mentales tengan algún papel en la explicación del comportamiento",
      "En creer que los estados cerebrales son reducibles a procesos conductuales observables"]),

    ("El reduccionismo como estrategia para resolver el problema mente-cuerpo propone:",
     ["Eliminar el vocabulario mental del lenguaje científico de la psicología",
      "Traducir las atribuciones de estados mentales a propiedades físicas o conductuales observables",
      "Aceptar el dualismo cartesiano como base filosófica de la psicología clínica",
      "Concebir la mente como una propiedad emergente irreducible al cerebro"]),

    ("¿Cuál es el principal problema del reduccionismo identificado desde el análisis de la conducta?",
     ["Que niega definitivamente la existencia de estados mentales privados",
      "Que hace imposible la intervención terapéutica eficaz en contextos naturales",
      "El problema de la realización múltiple: la misma conducta puede correlacionar con distintos estados cerebrales",
      "Que confunde las explicaciones normativas con las nomológicas"]),

    ("¿Qué propone Wittgenstein (1953) sobre el significado de las expresiones lingüísticas?",
     ["Que el significado depende de su capacidad para representar hechos del mundo",
      "Que el significado depende del uso de las expresiones en distintos juegos del lenguaje sociales",
      "Que el lenguaje mental es completamente privado e inaccesible a terceros por definición",
      "Que las atribuciones de estados mentales son pseudoproposiciones sin significado alguno"]),

    ("Cuando decimos 'Homero cree que los donuts son deliciosos', desde el anti-descriptivismo esto significa:",
     ["Que existe un estado neural real en el cerebro de Homero que representa los donuts",
      "Que Homero tiene activación dopaminérgica elevada ante la presencia de donuts",
      "Que podemos esperar de Homero ciertos cursos de acción: elegir donuts, buscarlos, hablar de ellos",
      "Que Homero ha sido reforzado positivamente con donuts a lo largo de su historia de aprendizaje"]),

    ("¿Cuál es la diferencia fundamental entre una explicación normativa y una nomológica?",
     ["Las normativas son más precisas científicamente que las nomológicas",
      "Las normativas explican causas físicas; las nomológicas explican razones e intenciones",
      "Las normativas establecen razones y evalúan racionalidad; las nomológicas establecen causas materiales",
      "No hay diferencia relevante; ambas tipos son equivalentes en psicología clínica"]),

    ("Desde el anti-descriptivismo, ¿qué evidencia determina si Homero realmente 'valora' su familia?",
     ["Lo que Homero dice explícitamente que valora cuando se le pregunta",
      "Lo que Homero se dice a sí mismo internamente (conducta verbal encubierta)",
      "Lo que Homero hace (comportamiento manifiesto y encubierto) en relación con su familia",
      "El diagnóstico clínico del terapeuta sobre la estructura de su personalidad"]),

    ("¿Por qué la neurociencia es insuficiente para explicar completamente la conducta?",
     ["Porque la neurociencia no es una disciplina científica empírica rigurosa",
      "Porque los procesos cerebrales no tienen ninguna influencia sobre el comportamiento",
      "Por razones epistemológicas (realización múltiple) y pragmáticas (no permite intervención conductual directa)",
      "Porque las explicaciones neurocientíficas solo aplican a animales, no a seres humanos"]),

    ("¿Qué son los 'nexos a distancia o temporales' que caracterizan al análisis de la conducta?",
     ["Las relaciones causales físico-contiguas entre neuronas que posibilitan la conducta",
      "Las relaciones físico-espaciales entre el organismo y los objetos de su entorno",
      "Las relaciones de contingencia temporal entre estímulos y respuestas, sin requerir contigüidad física",
      "Los vínculos terapéuticos entre terapeuta y paciente a lo largo del proceso clínico"]),

    ("El análisis de la conducta comparte con la teoría evolutiva de Darwin:",
     ["El interés por el sustrato genético y hereditario de la conducta compleja",
      "Un modelo causal basado en la selección por consecuencias a través del tiempo (nexos temporales)",
      "La idea de que la conducta humana es en su mayoría producto de instintos innatos",
      "El nivel de análisis subagencial centrado en procesos biológicos intraorganísmicos"]),

    ("Según el conductismo radical, ¿en qué difiere la conducta encubierta de la manifiesta?",
     ["La conducta encubierta no puede ser objeto de estudio científico empírico",
      "La conducta encubierta tiene naturaleza mental; la manifiesta, naturaleza física",
      "Solo difieren en el número de personas que pueden reportarlas, no en su naturaleza ontológica",
      "La encubierta es causada por estados cerebrales; la manifiesta, por contingencias ambientales"]),

    ("Cuando el psicólogo dice 'Homero come donuts porque tiene baja autoestima y falta de voluntad', ¿qué tipo de explicación es?",
     ["Una explicación nomológica basada en el análisis funcional de la conducta",
      "Una explicación normativa que evalúa la conducta según estándares, no establece causas",
      "Una explicación subagencial basada en neurociencia cognitiva aplicada",
      "Una hipótesis de mantenimiento válida y operativa para guiar el tratamiento"]),

    ("¿Cuál es la afirmación correcta sobre la relación entre neurociencia y análisis de la conducta?",
     ["El análisis de la conducta es un subconjunto de la neurociencia aplicada",
      "La neurociencia podrá reemplazar completamente al AF cuando mejore la tecnología",
      "Son niveles de análisis distintos y complementarios pero no reducibles entre sí",
      "El AF niega que el cerebro tenga cualquier influencia relevante sobre el comportamiento"]),

    ("¿Qué implica el 'principio de suficiencia neural' que el conductismo radical rechaza?",
     ["Que para cada tipo de conducta existe una disposición cerebral específica e invariable",
      "Que el cerebro es suficiente para sustentar las funciones vitales del organismo",
      "Que las leyes del aprendizaje son suficientes para explicar toda conducta animal",
      "Que la suficiencia del AF para explicar la conducta hace innecesaria la neurociencia"]),

    ("Desde el AF, ¿cómo se conceptualiza el pensamiento como 'habla interna'?",
     ["Como un proceso cognitivo interno que dirige y causa el comportamiento observable posterior",
      "Como conducta verbal encubierta: mismo tipo de respuesta que el habla manifiesta pero silenciosa",
      "Como un estado mental privado que solo puede conocerse mediante introspección directa",
      "Como activación secuencial de neuronas en el córtex prefrontal durante la planificación"]),

    ("La 'folk psychology' es:",
     ["Una teoría psicológica científica alternativa al conductismo radical",
      "El conjunto de explicaciones cotidianas del comportamiento en términos de creencias, deseos e intenciones",
      "El análisis funcional aplicado informalmente a situaciones cotidianas no clínicas",
      "La psicología cultural característica de comunidades rurales latinoamericanas"]),

    ("¿Qué significa que la conducta es 'la raíz de lo psicológico' para el conductismo radical?",
     ["Que toda conducta tiene una base genética y neurológica determinante",
      "Que los estados mentales son la causa última de todas las conductas del organismo",
      "Que lo que llamamos estados psicológicos se expresa y estudia a través de la conducta, manifiesta y encubierta",
      "Que la psicología científica debe limitarse al estudio de conductas observables por terceros"]),

    ("Carlos dice 'sé que debo hacer ejercicio' pero nunca lo hace. Desde el anti-descriptivismo, ¿qué se concluye?",
     ["Carlos tiene un esquema cognitivo disfuncional que bloquea la conducta de hacer ejercicio",
      "Su autoatribución mental ('sé que debo') no coincide con lo que su comportamiento expresa",
      "Carlos carece de motivación intrínseca suficiente para el ejercicio físico regular",
      "Hay un déficit neurológico que impide la planificación y ejecución de la conducta motora"]),
]

respuestas_U2 = ['B','C','B','B','C','B','C','C','C','C','C','B','C','B','C','A','B','B','C','B']

make_quiz(
    r'C:\Users\MI PC\psicoeduca\materiales\cuestionario-unidad2.docx',
    'Cuestionario — Bases Filosóficas del Análisis de la Conducta (Unidad 2)',
    'Seleccioná la única opción correcta para cada pregunta. '
    'Cada pregunta vale 0,5 puntos. Total: 10 puntos.',
    preguntas_U2,
    respuestas_U2
)
print('3/3 OK: cuestionario-unidad2.docx')
print('')
print('Todos los archivos guardados en: materiales/')
