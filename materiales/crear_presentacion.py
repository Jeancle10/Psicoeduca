"""
Presentacion: Análisis Funcional de la Conducta
Marca: PsicoEduca — identidad Rebranding
Autor: Lic. Jean Clemotte | @Psico_Educa20
Caso guia: Marge Simpson — Asuncion, Paraguay
"""
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
from pptx.oxml.ns import qn

# ── COLORES ─────────────────────────────────────────────
DARK_BG = RGBColor(0x1E, 0x3A, 0x5F)
CREMA   = RGBColor(0xF2, 0xED, 0xE4)
NAVY    = RGBColor(0x2B, 0x5E, 0xA7)
SKY     = RGBColor(0x4A, 0x9F, 0xE0)
ORANGE  = RGBColor(0xE8, 0xA8, 0x35)
GREEN_W = RGBColor(0x4A, 0xBF, 0xB0)
CREAM_T = RGBColor(0xF5, 0xF0, 0xDC)
DARK_T  = RGBColor(0x1E, 0x3A, 0x5F)
BOX_BG  = RGBColor(0xE4, 0xEE, 0xF8)
RED     = RGBColor(0xCC, 0x33, 0x33)
GREEN_OK= RGBColor(0x22, 0x88, 0x44)

SERIF = 'Georgia'
SANS  = 'Calibri'

W  = Cm(33.87)
H  = Cm(19.05)
ML = Cm(1.8)
CW = Cm(30.27)

LOGO_CREMA = r'C:\Users\MI PC\Desktop\PsicoEduca\Marketing\Rebranding\claro sin fondo.png'
LOGO_DARK  = r'C:\Users\MI PC\Desktop\PsicoEduca\Marketing\Rebranding\oscuro sin fondo.png'
PHOTO      = r'C:\Users\MI PC\Desktop\PsicoEduca\Marketing\BOOK\sentado.JPG'

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]


# ── HELPERS ─────────────────────────────────────────────

def bg(slide, color=CREMA):
    f = slide.background.fill
    f.solid(); f.fore_color.rgb = color


def hbar(slide, x, y, w, color=ORANGE, thick=Pt(3)):
    s = slide.shapes.add_shape(1, x, y, w, thick)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()


def txt(slide, text, x, y, w, h,
        font=SANS, size=16, bold=False, italic=False,
        color=DARK_T, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb


def paras(slide, lines, x, y, w, h,
          font=SANS, size=14, color=DARK_T, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if isinstance(ln, dict):
            r = p.add_run()
            r.text = ln.get('t', '')
            r.font.name   = ln.get('font', font)
            r.font.size   = Pt(ln.get('size', size))
            r.font.bold   = ln.get('bold', False)
            r.font.italic = ln.get('italic', False)
            r.font.color.rgb = ln.get('color', color)
        else:
            r = p.add_run(); r.text = str(ln)
            r.font.name = font; r.font.size = Pt(size)
            r.font.color.rgb = color
    return tb


def logo(slide, dark=False):
    path = LOGO_DARK if dark else LOGO_CREMA
    try:
        slide.shapes.add_picture(path, W - Cm(4.5), Cm(0.3), Cm(4.0), Cm(2.2))
    except Exception:
        pass


def footer(slide, dark=False):
    c = CREAM_T if dark else RGBColor(0x88, 0x88, 0x88)
    txt(slide, 'Lic. Jean Clemotte  |  @Psico_Educa20',
        ML, H - Cm(1.1), Cm(22), Cm(0.8),
        SANS, 10, italic=True, color=c)


def header_mod(slide, label):
    txt(slide, label.upper(), ML, Cm(0.5), Cm(22), Cm(0.7),
        SANS, 9, bold=True, color=NAVY)
    hbar(slide, ML, Cm(1.3), Cm(22), NAVY, Pt(1))


def box(slide, text, x, y, w, h,
        fill=BOX_BG, border=NAVY,
        font=SANS, size=14, bold=False, color=DARK_T,
        align=PP_ALIGN.LEFT):
    s = slide.shapes.add_shape(5, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = border; s.line.width = Pt(1)
    tf = s.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = color
    return s


def tbl(slide, data, x, y, w, h,
        hdr_bg=DARK_BG, hdr_fg=CREAM_T,
        odd=CREMA, even=RGBColor(0xE0, 0xDB, 0xD0),
        fs=12, center=False):
    rows = len(data); cols = max(len(r) for r in data)
    t = slide.shapes.add_table(rows, cols, x, y, w, h).table
    al = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
    for ri, row in enumerate(data):
        for ci in range(cols):
            val = row[ci] if ci < len(row) else ''
            cell = t.cell(ri, ci); cell.text = str(val)
            tf = cell.text_frame; tf.word_wrap = True
            for para in tf.paragraphs:
                para.alignment = al
                for run in para.runs:
                    run.font.name = SANS; run.font.size = Pt(fs)
                    run.font.bold = (ri == 0)
                    run.font.color.rgb = hdr_fg if ri == 0 else DARK_T
            tcPr = cell._tc.get_or_add_tcPr()
            for old in tcPr.findall(qn('a:solidFill')): tcPr.remove(old)
            sf  = etree.SubElement(tcPr, qn('a:solidFill'))
            clr = etree.SubElement(sf, qn('a:srgbClr'))
            if ri == 0:   clr.set('val', '{:02X}{:02X}{:02X}'.format(*hdr_bg))
            elif ri%2==1: clr.set('val', '{:02X}{:02X}{:02X}'.format(*odd))
            else:         clr.set('val', '{:02X}{:02X}{:02X}'.format(*even))
    return t


def bloque(num, titulo, sub=''):
    sl = prs.slides.add_slide(blank); bg(sl, DARK_BG)
    # Número grande naranja
    txt(sl, num, ML, Cm(3.0), Cm(5), Cm(6),
        SERIF, 96, bold=True, color=ORANGE)
    hbar(sl, ML + Cm(5.8), Cm(6.5), Cm(20), ORANGE, Pt(3))
    txt(sl, titulo, ML + Cm(5.8), Cm(7.2), CW - Cm(7), Cm(5),
        SERIF, 38, bold=True, color=CREAM_T)
    if sub:
        txt(sl, sub, ML + Cm(5.8), Cm(12.5), CW - Cm(7), Cm(2),
            SANS, 16, color=SKY)
    logo(sl, dark=True); footer(sl, dark=True)
    return sl


# ════════════════════════════════════════════════════════
# S1 — PORTADA  (DARK)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, DARK_BG)

txt(sl, 'Análisis\nFuncional\nde la Conducta',
    ML, Cm(2.0), Cm(22), Cm(11),
    SERIF, 52, bold=True, color=CREAM_T)

hbar(sl, ML, Cm(13.2), Cm(20), ORANGE, Pt(3))

txt(sl, 'Cómo entender el comportamiento más allá del diagnóstico',
    ML, Cm(13.8), Cm(25), Cm(1.4), SANS, 16, color=SKY)
txt(sl, 'Presentación por Lic. Jean Clemotte',
    ML, Cm(15.4), Cm(22), Cm(1.2),
    SERIF, 14, italic=True, color=RGBColor(0xAA, 0xBB, 0xCC))

# Elemento visual: texto "AF" enorme semitransparente derecha
txt(sl, 'AF', W - Cm(11), Cm(1.5), Cm(10), Cm(14),
    SERIF, 160, bold=True, color=RGBColor(0x2B, 0x4A, 0x6E),
    align=PP_ALIGN.CENTER)

logo(sl, dark=True); footer(sl, dark=True)


# ════════════════════════════════════════════════════════
# S2 — GANCHO: MARGE  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, CREMA)
header_mod(sl, 'Análisis Funcional de la Conducta')

txt(sl, 'Marge, 39 años — Asunción, Paraguay',
    ML, Cm(1.7), CW - Cm(5), Cm(2.0),
    SERIF, 30, bold=True, color=DARK_T)

paras(sl, [
    {'t': '"Homero convenció a Marge de ir a consulta.', 'size': 16, 'italic': True},
    {'t': ' La asaltaron en el Super Stock de Villa Morra.', 'size': 16, 'italic': True},
    {'t': ' Desde entonces evita salir sola a cualquier lugar concurrido."', 'size': 16, 'italic': True},
    ' ',
    {'t': 'Diagnóstico: Trastorno de pánico con agorafobia.', 'size': 15, 'bold': True, 'color': NAVY},
], ML, Cm(4.2), Cm(19), Cm(7), SANS, 16, DARK_T)

# Visual: "¿Eso explica?" + NO en naranja
txt(sl, '¿Eso explica\npor qué no\npuede salir?',
    Cm(22), Cm(3.8), Cm(10), Cm(4.5),
    SERIF, 20, color=NAVY, align=PP_ALIGN.CENTER)
txt(sl, 'NO.',
    Cm(22), Cm(8.8), Cm(10), Cm(4),
    SERIF, 72, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

hbar(sl, ML, Cm(13.5), CW - Cm(1), ORANGE)
txt(sl, 'Solo le pone nombre. Vamos a aprender a responder la pregunta real.',
    ML, Cm(13.9), CW - Cm(1), Cm(1.2),
    SANS, 15, italic=True, color=NAVY)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S3 — ÍNDICE  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, CREMA)
header_mod(sl, 'Análisis Funcional de la Conducta')

txt(sl, '¿Qué vamos a ver?',
    ML, Cm(1.7), Cm(12), Cm(2.0),
    SERIF, 34, bold=True, color=DARK_T)

items = [
    ('01', 'Fundamentos del AF', 'Qué es  •  4 supuestos  •  Funcionalidad vs Morfología'),
    ('02', 'Las leyes del aprendizaje', 'CC  •  CO  •  Ed  •  Variables disposicionales'),
    ('03', 'El procedimiento paso a paso', 'Describir  •  Antecedentes/consecuentes  •  Hipótesis funcional'),
]
for i, (num, tit, desc) in enumerate(items):
    y = Cm(4.5) + i * Cm(4.0)
    s = sl.shapes.add_shape(5, ML, y, Cm(2.3), Cm(1.9))
    s.fill.solid(); s.fill.fore_color.rgb = DARK_BG; s.line.fill.background()
    tf = s.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num; r.font.name = SERIF
    r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = ORANGE
    txt(sl, tit, ML + Cm(3.0), y, CW - Cm(4), Cm(1.1),
        SANS, 17, bold=True, color=DARK_T)
    txt(sl, desc, ML + Cm(3.0), y + Cm(1.2), CW - Cm(4), Cm(0.9),
        SANS, 13, color=NAVY)
    if i < 2:
        hbar(sl, ML, y + Cm(2.3), CW - Cm(1),
             RGBColor(0xCC, 0xC8, 0xBF), Pt(1))

txt(sl, 'Caso guía: Marge Simpson  —  Super Stock, Villa Morra',
    ML, Cm(16.5), CW - Cm(1), Cm(1.0),
    SANS, 13, italic=True, color=NAVY)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S4 — PORTADA BLOQUE 1  (DARK)
# ════════════════════════════════════════════════════════
bloque('01', '¿Qué es el\nAnálisis Funcional?',
       'Definición  •  4 supuestos  •  Funcionalidad vs Morfología')


# ════════════════════════════════════════════════════════
# S5 — DIAGNÓSTICO VS AF  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, CREMA)
header_mod(sl, 'Bloque 01 — ¿Qué es el AF?')

txt(sl, 'El diagnóstico describe. El AF explica.',
    ML, Cm(1.7), CW - Cm(5), Cm(1.8),
    SERIF, 30, bold=True, color=DARK_T)

data5 = [
    ['Diagnóstico', 'El Análisis Funcional pregunta...'],
    ['"Trastorno de pánico"',  '¿Qué hace Marge exactamente?'],
    ['"Depresión mayor"',      '¿Cuándo ocurre? ¿Cuándo NO ocurre?'],
    ['"Fobia social"',         '¿Qué pasa después de que lo hace?'],
]
tbl(sl, data5, ML, Cm(4.2), CW - Cm(5), Cm(8.0), fs=16)

# Visual: flecha derecha con frase clave
txt(sl, '→', Cm(26.5), Cm(7.5), Cm(3), Cm(3),
    SANS, 60, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
txt(sl, 'El AF busca\nrelaciones entre\nla conducta\ny su contexto.',
    Cm(26.5), Cm(10.5), Cm(6), Cm(5),
    SANS, 14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

hbar(sl, ML, Cm(13.8), CW - Cm(1), ORANGE)
txt(sl, 'El diagnóstico es necesario para la administración clínica — pero no explica las CAUSAS ni permite intervenir.',
    ML, Cm(14.2), CW - Cm(1), Cm(1.5), SANS, 14, italic=True, color=NAVY)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S6 — FUNCIONALIDAD VS MORFOLOGÍA  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, CREMA)
header_mod(sl, 'Bloque 01 — ¿Qué es el AF?')

txt(sl, 'Lo que importa no es cómo se ve, sino para qué sirve.',
    ML, Cm(1.7), CW - Cm(5), Cm(1.8),
    SERIF, 28, bold=True, color=DARK_T)

txt(sl, 'Bart quiere evitar ir a la escuela...',
    ML, Cm(4.0), CW, Cm(1.0), SANS, 16, bold=True, color=NAVY)

behaviors = ['Dice que\nle duele\nla panza', 'Esconde\nla mochila',
             'Hace una\nrabieta', 'Llora', 'Hace el\ntrabajo mal']
bw = CW / 5 - Cm(0.4)
for i, b in enumerate(behaviors):
    x = ML + i * (bw + Cm(0.4))
    s = sl.shapes.add_shape(5, x, Cm(5.5), bw, Cm(3.0))
    s.fill.solid(); s.fill.fore_color.rgb = BOX_BG
    s.line.color.rgb = NAVY; s.line.width = Pt(1)
    tf = s.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = b
    r.font.name = SANS; r.font.size = Pt(14); r.font.color.rgb = DARK_T

# Visual: llaves que convergen → función única
txt(sl, '5 morfologías distintas', ML, Cm(9.5), CW/2, Cm(1.0),
    SANS, 16, color=DARK_T)
txt(sl, '→', CW/2 + ML, Cm(9.3), Cm(2.5), Cm(1.5),
    SANS, 36, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
box(sl, '1 función:\nescapar de la escuela',
    CW/2 + ML + Cm(3), Cm(8.8), Cm(9), Cm(2.5),
    fill=DARK_BG, border=DARK_BG,
    font=SANS, size=16, bold=True, color=CREAM_T, align=PP_ALIGN.CENTER)

hbar(sl, ML, Cm(12.8), CW - Cm(1), GREEN_W)
txt(sl, '"El AF no pregunta cómo se ve la conducta, sino qué función cumple en ese contexto."',
    ML, Cm(13.2), CW - Cm(1), Cm(1.8),
    SERIF, 16, italic=True, color=NAVY)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S7 — 4 SUPUESTOS  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, CREMA)
header_mod(sl, 'Bloque 01 — ¿Qué es el AF?')

txt(sl, 'Antes de empezar: 4 supuestos',
    ML, Cm(1.7), CW - Cm(5), Cm(1.8),
    SERIF, 30, bold=True, color=DARK_T)

supuestos = [
    ('01', ORANGE,  'Analiza CUALQUIER conducta',
     'Incluyendo pensamientos, emociones y conductas encubiertas'),
    ('02', GREEN_W, 'Las conductas son ADAPTATIVAS',
     'A corto plazo. El problema es que dejan de funcionar a largo plazo'),
    ('03', ORANGE,  'Las leyes son UNIVERSALES',
     'Aplican igual en Paraguay, España o cualquier parte del mundo'),
    ('04', GREEN_W, 'Considera MUCHAS variables',
     'Biológicas, psicológicas, históricas y contextuales'),
]
hw = CW / 2 - Cm(0.6)
for i, (num, accent, tit, desc) in enumerate(supuestos):
    col = i % 2; row = i // 2
    x = ML + col * (hw + Cm(0.8))
    y = Cm(4.2) + row * Cm(5.8)
    s = sl.shapes.add_shape(1, x, y, hw, Cm(5.3))
    s.fill.solid(); s.fill.fore_color.rgb = BOX_BG
    s.line.color.rgb = accent; s.line.width = Pt(2)
    # Número
    txt(sl, num, x + Cm(0.4), y + Cm(0.3), Cm(2), Cm(1.2),
        SERIF, 22, bold=True, color=accent)
    txt(sl, tit, x + Cm(0.4), y + Cm(1.6), hw - Cm(0.8), Cm(1.2),
        SANS, 14, bold=True, color=DARK_T)
    txt(sl, desc, x + Cm(0.4), y + Cm(2.9), hw - Cm(0.8), Cm(2.0),
        SANS, 13, color=NAVY)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S8 — PORTADA BLOQUE 2  (DARK)
# ════════════════════════════════════════════════════════
bloque('02', 'Las leyes\ndel aprendizaje',
       'Habituación  •  CC  •  CO  •  Ed  •  Variables')


# ════════════════════════════════════════════════════════
# S9 — HABITUACIÓN Y SENSIBILIZACIÓN  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, CREMA)
header_mod(sl, 'Bloque 02 — Las leyes del aprendizaje')

txt(sl, 'Habituación y Sensibilización',
    ML, Cm(1.7), CW - Cm(5), Cm(1.8),
    SERIF, 30, bold=True, color=DARK_T)

hw9 = CW / 2 - Cm(0.6)
# Header HABITUACIÓN
s1 = sl.shapes.add_shape(5, ML, Cm(4.2), hw9, Cm(1.4))
s1.fill.solid(); s1.fill.fore_color.rgb = DARK_BG; s1.line.fill.background()
tf = s1.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = 'HABITUACIÓN'
r.font.name = SANS; r.font.size = Pt(15); r.font.bold = True
r.font.color.rgb = ORANGE

paras(sl, [
    {'t': 'La respuesta BAJA con la exposición repetida.', 'bold': True, 'size': 15},
    {'t': '→ El estímulo pierde fuerza', 'size': 14, 'color': NAVY},
    ' ',
    {'t': 'Ej: Dejar de oír el ruido del AC de la oficina', 'size': 14, 'italic': True},
    {'t': 'Ej: No notar el olor a tereré en la sala', 'size': 14, 'italic': True},
], ML, Cm(5.9), hw9, Cm(8.5), SANS, 14, DARK_T)

# Header SENSIBILIZACIÓN
col2 = ML + hw9 + Cm(0.8)
s2 = sl.shapes.add_shape(5, col2, Cm(4.2), hw9, Cm(1.4))
s2.fill.solid(); s2.fill.fore_color.rgb = NAVY; s2.line.fill.background()
tf2 = s2.text_frame; p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = 'SENSIBILIZACIÓN'
r2.font.name = SANS; r2.font.size = Pt(15); r2.font.bold = True
r2.font.color.rgb = GREEN_W

paras(sl, [
    {'t': 'La respuesta SUBE con la exposición repetida.', 'bold': True, 'size': 15},
    {'t': '→ El estímulo gana fuerza', 'size': 14, 'color': NAVY},
    ' ',
    {'t': 'Ej: Golpecitos que se sienten cada vez más fuertes', 'size': 14, 'italic': True},
    {'t': 'Ej: Ruido que parece cada vez más intenso', 'size': 14, 'italic': True},
], col2, Cm(5.9), hw9, Cm(8.5), SANS, 14, DARK_T)

hbar(sl, ML, Cm(15.3), CW - Cm(1), ORANGE)
txt(sl, '⚠  Habituación ≠ "acostumbrarse"  (eso es condicionamiento operante)',
    ML, Cm(15.7), CW - Cm(1), Cm(1.1),
    SANS, 14, bold=True, color=NAVY)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S10 — CONDICIONAMIENTO CLÁSICO  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, CREMA)
header_mod(sl, 'Bloque 02 — Las leyes del aprendizaje')

txt(sl, 'Condicionamiento Clásico: aprender a reaccionar',
    ML, Cm(1.7), CW - Cm(5), Cm(1.8),
    SERIF, 26, bold=True, color=DARK_T)

data10 = [
    ['Momento', 'Estímulo', '', 'Respuesta'],
    ['ANTES', 'Super Stock (neutro)', '→', 'Sin respuesta de ansiedad'],
    ['', 'Asalto en el Super Stock (EI)', '→', 'Miedo + taquicardia (RI)'],
    ['DURANTE', 'Super Stock + Asalto (EI)', '→', 'Miedo + taquicardia (RI)'],
    ['', '[Una sola vez — experiencia muy intensa]', '', ''],
    ['DESPUÉS', 'Super Stock (EC)', '→', 'Miedo + taquicardia (RC)'],
    ['', '↓ generalización: cualquier supermercado (EC)', '→', 'Miedo (RC)'],
]
tbl(sl, data10, ML, Cm(4.0), Cm(24), Cm(11.0), fs=13)

# Visual: nota al costado
box(sl, 'EI = estímulo biológicamente amenazante.\nEl pánico es la RI,\nno el EI.',
    Cm(26.5), Cm(5.5), Cm(6.5), Cm(4.0),
    fill=RGBColor(0xFF, 0xF0, 0xD8), border=ORANGE,
    font=SANS, size=13, color=DARK_T)

txt(sl, 'Las conductas respondientes ocurren POR algo — no dependen de lo que haga el organismo después.',
    ML, Cm(15.8), CW - Cm(1), Cm(1.3),
    SANS, 14, italic=True, color=NAVY)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S11 — TRIPLE CONTINGENCIA  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, CREMA)
header_mod(sl, 'Bloque 02 — Las leyes del aprendizaje')

txt(sl, 'Triple Contingencia',
    ML, Cm(1.7), CW - Cm(5), Cm(1.8),
    SERIF, 32, bold=True, color=DARK_T)

bw = Cm(8.8); bh = Cm(7.0); by = Cm(4.0); gap = Cm(1.1)

box(sl, 'ANTECEDENTE\n(Señal)\n\nHomero dice:\n"¿Vamos al\nSuper Stock?"\n\n(Ed)',
    ML, by, bw, bh,
    fill=BOX_BG, border=DARK_BG, size=15)

txt(sl, '▶', ML + bw + Cm(0.1), by + bh/2 - Cm(0.7), gap, Cm(1.5),
    SANS, 28, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

box(sl, 'RESPUESTA\n(Lo que hace)\n\n"Me duele\nla cabeza,\nno puedo ir."\n\n(Operante)',
    ML + bw + gap, by, bw, bh,
    fill=BOX_BG, border=DARK_BG, size=15)

txt(sl, '▶', ML + bw*2 + gap + Cm(0.1), by + bh/2 - Cm(0.7), gap, Cm(1.5),
    SANS, 28, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

box(sl, 'CONSECUENTE\n(Lo que pasa)\n\nHomero va solo.\nMarge en casa.\nDesaparece\nla ansiedad.\n(RN)',
    ML + bw*2 + gap*2, by, bw, bh,
    fill=BOX_BG, border=DARK_BG, size=15)

hbar(sl, ML, Cm(12.3), CW - Cm(1), GREEN_W)
txt(sl, 'Las conductas operantes ocurren PARA algo — la historia de consecuencias determina si se repiten.',
    ML, Cm(12.7), CW - Cm(1), Cm(1.3),
    SANS, 14, italic=True, color=NAVY)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S12 — ESTÍMULO DISCRIMINATIVO  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, CREMA)
header_mod(sl, 'Bloque 02 — Las leyes del aprendizaje')

txt(sl, 'Antecedente vs Estímulo Discriminativo (Ed)',
    ML, Cm(1.7), CW - Cm(5), Cm(1.8),
    SERIF, 28, bold=True, color=DARK_T)

hw12 = CW / 2 - Cm(0.6)
box(sl, 'ANTECEDENTE', ML, Cm(4.2), hw12, Cm(1.3),
    fill=BOX_BG, border=NAVY, font=SANS, size=16, bold=True, color=DARK_T,
    align=PP_ALIGN.CENTER)
paras(sl, [
    {'t': 'Concepto DESCRIPTIVO', 'bold': True, 'size': 15, 'color': NAVY},
    'Solo indica relación temporal.',
    'El estímulo está ANTES de la conducta.',
    ' ',
    {'t': 'No dice nada sobre si hay reforzador disponible.', 'size': 13, 'color': NAVY, 'italic': True},
], ML, Cm(5.8), hw12, Cm(8), SANS, 14, DARK_T)

col2 = ML + hw12 + Cm(0.8)
box(sl, 'ESTÍMULO DISCRIMINATIVO (Ed)', col2, Cm(4.2), hw12, Cm(1.3),
    fill=DARK_BG, border=DARK_BG, font=SANS, size=14, bold=True, color=CREAM_T,
    align=PP_ALIGN.CENTER)
paras(sl, [
    {'t': 'Concepto FUNCIONAL', 'bold': True, 'size': 15, 'color': ORANGE},
    'Señala que si emito esta respuesta,',
    {'t': 'hay probabilidad de obtener el reforzador.', 'bold': True},
    ' ',
    {'t': 'Ed+ — Homero de buen humor → pedir permiso', 'size': 13, 'italic': True, 'color': GREEN_W},
    {'t': 'Ed− — Homero de mal humor → mejor no pedir', 'size': 13, 'italic': True, 'color': ORANGE},
], col2, Cm(5.8), hw12, Cm(9), SANS, 14, DARK_T)

txt(sl, 'No todo lo que ocurre antes es un Ed. El Ed se establece por la historia de aprendizaje.',
    ML, Cm(15.6), CW - Cm(1), Cm(1.2),
    SANS, 14, bold=True, italic=True, color=NAVY)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S13 — 4 PROCEDIMIENTOS OPERANTES  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, CREMA)
header_mod(sl, 'Bloque 02 — Las leyes del aprendizaje')

txt(sl, '4 procedimientos operantes básicos',
    ML, Cm(1.7), CW - Cm(5), Cm(1.8),
    SERIF, 30, bold=True, color=DARK_T)

data13 = [
    ['', 'APARECE algo después', 'DESAPARECE algo después'],
    ['La R\nSUBE',
     'Reforzamiento POSITIVO\nMarge apaga la TV → Bart grita → Marge la vuelve a prender (cede)',
     'Reforzamiento NEGATIVO\nMarge evita el Super Stock → desaparece el miedo → sigue evitando'],
    ['La R\nBAJA',
     'Castigo POSITIVO\nBart hace broma → Homero grita → Bart deja de hacerlas',
     'Castigo NEGATIVO\nBart llega tarde → le quitan el skate → llega a horario'],
]
tbl(sl, data13, ML, Cm(4.2), CW - Cm(1), Cm(11.0), fs=14, center=False)

hbar(sl, ML, Cm(16.2), CW - Cm(1), ORANGE)
txt(sl, '⚠  Positivo/Negativo = aparece/desaparece. No significa bueno/malo.',
    ML, Cm(16.6), CW - Cm(1), Cm(1.0),
    SANS, 14, bold=True, color=NAVY)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S14 — VARIABLES  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, CREMA)
header_mod(sl, 'Bloque 02 — Las leyes del aprendizaje')

txt(sl, 'Variables que alteran la contingencia',
    ML, Cm(1.7), CW - Cm(5), Cm(1.8),
    SERIF, 28, bold=True, color=DARK_T)

hw14 = CW / 2 - Cm(0.6)
s1h = sl.shapes.add_shape(5, ML, Cm(4.2), hw14, Cm(1.4))
s1h.fill.solid(); s1h.fill.fore_color.rgb = DARK_BG; s1h.line.fill.background()
tf = s1h.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = 'DISPOSICIONALES (estables)'
r.font.name = SANS; r.font.size = Pt(14); r.font.bold = True
r.font.color.rgb = ORANGE

paras(sl, [
    '• Biológicas: edad, enfermedades, medicación',
    '• Historia de aprendizaje',
    '• Repertorio conductual (habilidades/déficits)',
    {'t': '• Reglas internas: "si hago X, pasa Y"', 'color': ORANGE, 'bold': True},
    '• Condiciones del entorno físico y social',
], ML, Cm(5.9), hw14, Cm(8), SANS, 14, DARK_T)

col2 = ML + hw14 + Cm(0.8)
s2h = sl.shapes.add_shape(5, col2, Cm(4.2), hw14, Cm(1.4))
s2h.fill.solid(); s2h.fill.fore_color.rgb = NAVY; s2h.line.fill.background()
tf2 = s2h.text_frame; p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = 'MOTIVADORAS (cambian)'
r2.font.name = SANS; r2.font.size = Pt(14); r2.font.bold = True
r2.font.color.rgb = GREEN_W

paras(sl, [
    '• Privación / saciación',
    '• Estado emocional actual',
    '• Anticipaciones verbales',
    {'t': '• Cambios en el contexto inmediato', 'color': GREEN_W, 'bold': True},
], col2, Cm(5.9), hw14, Cm(6), SANS, 14, DARK_T)

txt(sl, 'No causan la conducta. Modifican la PROBABILIDAD de que ocurra.',
    ML, Cm(15.3), CW - Cm(1), Cm(1.2),
    SANS, 15, bold=True, italic=True, color=NAVY)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S15 — EJEMPLOS VARIABLES PARAGUAY  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, CREMA)
header_mod(sl, 'Bloque 02 — Las leyes del aprendizaje')

txt(sl, 'Ejemplos contextualizados — Marge en Asunción',
    ML, Cm(1.7), CW - Cm(5), Cm(1.8),
    SERIF, 28, bold=True, color=DARK_T)

data15 = [
    ['Variable', 'Tipo', 'Efecto sobre la conducta de Marge'],
    ['Asalto previo en el Super Stock',
     'Disposicional\n(historia)',
     'El Super Stock adquirió función de EC aversivo.\nAumenta la probabilidad de evitación.'],
    ['Regla: "Si salgo sola, me van a asaltar de nuevo"',
     'Disposicional\n(regla)',
     'Amplifica el valor del Ed. Aumenta la taquicardia ante cualquier supermercado.'],
    ['Estrés del hogar: Bart suspendió, Maggie no duerme',
     'Motivadora\n(estado emocional)',
     'Aumenta activación general. Cualquier estímulo produce más ansiedad.'],
    ['Homero siempre va solo al Super Stock si Marge no puede',
     'Motivadora\n(abolición)',
     'Reduce presión por salir. Refuerza la evitación por partida doble.'],
]
tbl(sl, data15, ML, Cm(4.0), CW - Cm(1), Cm(11.5), fs=12)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S16 — PORTADA BLOQUE 3  (DARK)
# ════════════════════════════════════════════════════════
bloque('03', 'El procedimiento\npaso a paso',
       'Describir la conducta  •  Antecedentes y consecuentes  •  Hipótesis funcional')


# ════════════════════════════════════════════════════════
# S17 — PASO 1  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, CREMA)
header_mod(sl, 'Bloque 03 — El procedimiento')

txt(sl, 'Paso 1 — ¿Qué hace exactamente?',
    ML, Cm(1.7), CW - Cm(5), Cm(1.8),
    SERIF, 30, bold=True, color=DARK_T)

dims = [
    ('01', 'Frecuencia', '¿Cuántas\nveces ocurre?'),
    ('02', 'Duración',   '¿Cuánto\ntiempo dura?'),
    ('03', 'Localización', '¿Cuándo?\n¿Dónde?\n¿Con quién?'),
]
bw17 = CW / 3 - Cm(0.5)
for i, (num, nom, preg) in enumerate(dims):
    x = ML + i * (bw17 + Cm(0.5))
    s = sl.shapes.add_shape(5, x, Cm(4.3), bw17, Cm(4.5))
    s.fill.solid(); s.fill.fore_color.rgb = BOX_BG
    s.line.color.rgb = NAVY; s.line.width = Pt(1)
    tf = s.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run(); r1.text = num + '\n'
    r1.font.name = SERIF; r1.font.size = Pt(28); r1.font.bold = True
    r1.font.color.rgb = ORANGE
    r2 = p.add_run(); r2.text = nom + '\n\n' + preg
    r2.font.name = SANS; r2.font.size = Pt(15); r2.font.color.rgb = DARK_T

txt(sl, '❌  "Marge tiene ansiedad"',
    ML, Cm(9.8), CW - Cm(1), Cm(1.1),
    SANS, 16, bold=True, color=RED)
txt(sl, '✓  "Marge evita supermercados, shoppings y eventos sociales desde hace 2 años en Asunción,\n   en cualquier lugar concurrido donde no esté Homero presente."',
    ML, Cm(11.2), CW - Cm(1), Cm(2.8), SANS, 15, color=GREEN_OK)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S18 — PASO 2  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, CREMA)
header_mod(sl, 'Bloque 03 — El procedimiento')

txt(sl, 'Paso 2 — ¿Qué pasa antes y después?',
    ML, Cm(1.7), CW - Cm(5), Cm(1.8),
    SERIF, 28, bold=True, color=DARK_T)

bw18 = Cm(9.0); bh18 = Cm(7.0); by18 = Cm(4.2); gap18 = Cm(1.2)

box(sl, 'ANTES\n(Antecedentes)\n\n¿Dónde está?\n¿Con quién?\n¿Qué piensa?\n¿Qué siente?',
    ML, by18, bw18, bh18, fill=BOX_BG, border=DARK_BG, size=15)

txt(sl, '▶', ML + bw18 + Cm(0.1), by18 + bh18/2 - Cm(0.7), gap18, Cm(1.5),
    SANS, 28, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

box(sl, 'CONDUCTA',
    ML + bw18 + gap18, by18 + bh18/2 - Cm(1.2), Cm(6.5), Cm(2.5),
    fill=DARK_BG, border=DARK_BG, font=SANS, size=18, bold=True, color=CREAM_T,
    align=PP_ALIGN.CENTER)

txt(sl, '▶', ML + bw18*2 + gap18 + Cm(0.3), by18 + bh18/2 - Cm(0.7), gap18, Cm(1.5),
    SANS, 28, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

box(sl, 'DESPUÉS\n(Consecuentes)\n\n¿Qué obtiene?\n¿Qué evita?\n¿Cómo se siente?\n¿Qué hacen los otros?',
    ML + bw18*2 + gap18*2 - Cm(0.2), by18, bw18, bh18,
    fill=BOX_BG, border=DARK_BG, size=15)

hbar(sl, ML, Cm(12.5), CW - Cm(1), ORANGE)
txt(sl, '⚠  No todo lo que ocurre antes/después tiene relación funcional con la conducta.',
    ML, Cm(12.9), CW - Cm(1), Cm(1.1),
    SANS, 14, bold=True, color=NAVY)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S19 — E-R VS E-R-C  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, CREMA)
header_mod(sl, 'Bloque 03 — El procedimiento')

txt(sl, '¿La conducta ocurre POR algo o PARA algo?',
    ML, Cm(1.7), CW - Cm(5), Cm(1.8),
    SERIF, 28, bold=True, color=DARK_T)

data19 = [
    ['E-R  (Respondiente / Clásica)', 'E-R-C  (Operante)'],
    ['Es una REACCIÓN ante el ambiente', 'Es una ACCIÓN sobre el ambiente'],
    ['No depende de las consecuencias', 'Las consecuencias la moldean'],
    ['Es automática — no se puede evitar', 'El sujeto puede modificarla'],
    ['Ocurre POR el estímulo antecedente', 'Ocurre PARA obtener algo'],
    ['Ej: taquicardia ante el Super Stock', 'Ej: decir "no puedo ir" y quedarse'],
]
tbl(sl, data19, ML, Cm(4.2), CW - Cm(1), Cm(9.5), fs=16, center=True)

hbar(sl, ML, Cm(14.5), CW - Cm(1), GREEN_W)
txt(sl, 'En clínica casi siempre aparecen las dos juntas: Marge tiene taquicardia (E-R)  y  evita el supermercado (E-R-C).',
    ML, Cm(14.9), CW - Cm(1), Cm(2.0),
    SANS, 15, italic=True, color=NAVY)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S20 — PASO 3  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, CREMA)
header_mod(sl, 'Bloque 03 — El procedimiento')

txt(sl, 'Paso 3 — ¿Qué facilita o dificulta que ocurra?',
    ML, Cm(1.7), CW - Cm(5), Cm(1.8),
    SERIF, 26, bold=True, color=DARK_T)

cuadros = [
    ('Biológicas',           '¿Medicación?\n¿Enfermedades?\n¿Duerme bien?'),
    ('Historia aprendizaje', '¿Cuándo empezó?\n¿Qué pasó?\n¿Experiencias similares?'),
    ('Reglas',               '"Si salgo, me asaltan"\n"No puedo controlar\nmi ansiedad"'),
    ('Contexto',             '¿Cuándo ocurre más?\n¿Cuándo NO ocurre?\n¿Qué lugares influyen?'),
]
hw20 = CW / 2 - Cm(0.6)
for i, (tit, cont) in enumerate(cuadros):
    col = i % 2; row = i // 2
    x = ML + col * (hw20 + Cm(0.8))
    y = Cm(4.2) + row * Cm(5.5)
    accent = ORANGE if col == 0 else GREEN_W
    txt(sl, tit, x, y, hw20, Cm(1.0),
        SANS, 15, bold=True, color=accent)
    hbar(sl, x, y + Cm(1.0), hw20, accent, Pt(2))
    txt(sl, cont, x, y + Cm(1.3), hw20, Cm(4.0),
        SANS, 14, color=DARK_T)

txt(sl, 'No listar por listar — especificar CÓMO cada variable afecta a la contingencia.',
    ML, Cm(16.2), CW - Cm(1), Cm(1.1),
    SANS, 14, italic=True, color=NAVY)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S21 — HIPÓTESIS FUNCIONAL  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, CREMA)
header_mod(sl, 'Bloque 03 — El procedimiento')

txt(sl, 'La hipótesis funcional: origen vs mantenimiento',
    ML, Cm(1.7), CW - Cm(5), Cm(1.8),
    SERIF, 26, bold=True, color=DARK_T)

hw21 = CW / 2 - Cm(0.6)
box(sl, 'Hipótesis de ORIGEN',
    ML, Cm(4.2), hw21, Cm(1.3),
    fill=BOX_BG, border=NAVY, font=SANS, size=16, bold=True, color=DARK_T,
    align=PP_ALIGN.CENTER)
paras(sl, [
    '¿Cómo se aprendió esta conducta?',
    ' ',
    {'t': '→ Puede ser desconocida', 'size': 13, 'color': NAVY},
    {'t': '→ No siempre necesaria para el tratamiento', 'size': 13, 'color': NAVY},
    ' ',
    {'t': '"Marge fue asaltada en el Super Stock.\nDesde ese día empezó a evitar salir."',
     'italic': True, 'size': 14},
], ML, Cm(5.8), hw21, Cm(10), SANS, 14, DARK_T)

col2 = ML + hw21 + Cm(0.8)
box(sl, 'Hipótesis de MANTENIMIENTO',
    col2, Cm(4.2), hw21, Cm(1.3),
    fill=DARK_BG, border=DARK_BG, font=SANS, size=14, bold=True, color=CREAM_T,
    align=PP_ALIGN.CENTER)
paras(sl, [
    '¿Qué la sostiene HOY?',
    ' ',
    {'t': '→ Esta es la que guía el tratamiento', 'bold': True, 'size': 13, 'color': ORANGE},
    {'t': '→ Se establece desde las contingencias actuales', 'size': 13},
    ' ',
    {'t': '"Marge evita porque eso reduce la taquicardia (RN). Homero resuelve las compras (RN adicional). La regla amplifica la respuesta."',
     'italic': True, 'size': 14},
], col2, Cm(5.8), hw21, Cm(10), SANS, 14, DARK_T)

txt(sl, 'El pasado no CAUSA la conducta actual. La sostienen las contingencias PRESENTES.',
    ML, Cm(16.3), CW - Cm(1), Cm(1.0),
    SANS, 15, bold=True, italic=True, color=NAVY)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S22 — VOLVEMOS A MARGE  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, CREMA)
header_mod(sl, 'Bloque 03 — El procedimiento')

txt(sl, 'Marge: del diagnóstico al análisis funcional',
    ML, Cm(1.7), CW - Cm(5), Cm(1.8),
    SERIF, 26, bold=True, color=DARK_T)

data22 = [
    ['Componente', 'Marge Simpson — Asunción, Paraguay'],
    ['Conducta (R)',        'Evitar supermercados, shoppings y lugares concurridos sin compañía'],
    ['Antecedente (Ed)',    'Homero dice "¿vamos al Super Stock?" / pensar en salir sola'],
    ['Consecuente (C)',     'Homero va solo → Marge en casa → desaparece ansiedad (RN)'],
    ['Resp. condicionada',  'Taquicardia + ahogo ante cualquier lugar concurrido (CC)'],
    ['Var. disposicional',  'Regla: "Si salgo sola, me van a asaltar de nuevo"'],
    ['Var. motivadora',     'Estrés del hogar: Bart + Homero + Maggie potencia la respuesta'],
]
tbl(sl, data22, ML, Cm(4.0), CW - Cm(1), Cm(10.5), fs=14)

hbar(sl, ML, Cm(15.5), CW - Cm(1), ORANGE)
txt(sl, 'La evitación se mantiene por DOS razones: produce alivio + Homero resuelve por ella.',
    ML, Cm(15.9), CW - Cm(1), Cm(1.2),
    SANS, 15, bold=True, color=NAVY)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S23 — IDEAS CLAVE  (DARK)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, DARK_BG)

txt(sl, 'Para llevarse hoy',
    ML, Cm(1.5), CW - Cm(5), Cm(1.8),
    SERIF, 36, bold=True, color=CREAM_T)

ideas = [
    ('01', ORANGE,  'Lo que importa no es cómo se ve la conducta (morfología),\nsino para qué sirve en ese contexto (función).'),
    ('02', GREEN_W, 'Las conductas se mantienen porque obtienen consecuencias.\nIdentificar cuáles es la tarea del psicólogo.'),
    ('03', ORANGE,  'El AF no es un diagnóstico.\nEs una hipótesis explicativa que guía la intervención.'),
]
for i, (num, accent, texto) in enumerate(ideas):
    y = Cm(4.5) + i * Cm(4.0)
    s = sl.shapes.add_shape(5, ML, y, Cm(2.5), Cm(3.0))
    s.fill.solid(); s.fill.fore_color.rgb = accent; s.line.fill.background()
    tf = s.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num
    r.font.name = SERIF; r.font.size = Pt(30); r.font.bold = True
    r.font.color.rgb = DARK_BG
    txt(sl, texto, ML + Cm(3.2), y + Cm(0.3), CW - Cm(7.5), Cm(2.8),
        SANS, 18, color=CREAM_T)

logo(sl, dark=True); footer(sl, dark=True)


# ════════════════════════════════════════════════════════
# S24 — CIERRE CON FOTO  (DARK)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
try:
    sl.shapes.add_picture(PHOTO, 0, 0, W, H)
except Exception:
    pass

overlay = sl.shapes.add_shape(1, 0, 0, W * 0.54, H)
overlay.fill.solid(); overlay.fill.fore_color.rgb = DARK_BG
overlay.line.fill.background()
xPr = overlay.fill._xPr; sf = xPr.solidFill
from pptx.oxml.ns import qn as _qn; from lxml import etree as _et
clr = sf.find(_qn('a:srgbClr'))
if clr is None:
    clr = _et.SubElement(sf, _qn('a:srgbClr')); clr.set('val', '1E3A5F')
alpha = _et.SubElement(clr, _qn('a:alpha')); alpha.set('val', '85000')

hbar(sl, ML, Cm(1.8), Cm(16), ORANGE, Pt(2))
hbar(sl, ML, Cm(17.0), Cm(16), ORANGE, Pt(2))

txt(sl, '@PSICO_EDUCA20', ML, Cm(0.5), Cm(16), Cm(0.9),
    SANS, 12, color=RGBColor(0xAA, 0xBB, 0xCC))

txt(sl, 'Gracias por\nsu Atención',
    ML, Cm(3.5), Cm(16), Cm(7.0),
    SERIF, 52, bold=True, italic=True, color=CREAM_T)

txt(sl, 'Análisis Funcional de la Conducta',
    ML, Cm(11.5), Cm(16), Cm(2.0), SANS, 17, color=SKY)

txt(sl, 'Lic. Jean Clemotte  |  PsicoEduca',
    ML, Cm(13.8), Cm(16), Cm(1.3), SANS, 15, italic=True, color=CREAM_T)

logo(sl, dark=True)


# ════════════════════════════════════════════════════════
# GUARDAR
# ════════════════════════════════════════════════════════
output = r'C:\Users\MI PC\psicoeduca\materiales\presentacion-AF.pptx'
prs.save(output)
print('Listo: ' + output)
print('24 slides | identidad Rebranding | Lic. Jean Clemotte')
