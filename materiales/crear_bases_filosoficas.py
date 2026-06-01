"""
Presentacion: Bases Filosoficas del Analisis de la Conducta
Marca: PsicoEduca — nueva identidad visual
Autora: Lic. Jean Clemotte | @Psico_Educa20
Fuente: Cap. 3 — Frojan et al.
"""
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
from pptx.oxml.ns import qn

# ── COLORES nueva identidad ──────────────────────────────
DARK_BG  = RGBColor(0x1E, 0x3A, 0x5F)   # azul marino oscuro (fondo dark)
CREMA    = RGBColor(0xF2, 0xED, 0xE4)   # crema suave (fondo light)
NAVY     = RGBColor(0x2B, 0x5E, 0xA7)   # azul PSE medio
SKY      = RGBColor(0x4A, 0x9F, 0xE0)   # azul cielo
ORANGE   = RGBColor(0xE8, 0xA8, 0x35)   # naranja acento
GREEN_W  = RGBColor(0x4A, 0xBF, 0xB0)   # verde agua acento
CREAM_T  = RGBColor(0xF5, 0xF0, 0xDC)   # texto sobre oscuro
DARK_T   = RGBColor(0x1E, 0x3A, 0x5F)   # texto sobre claro

SERIF = 'Georgia'
SANS  = 'Calibri'

W  = Cm(33.87)
H  = Cm(19.05)
ML = Cm(1.8)
CW = Cm(30.27)

# ── LOGOS (Rebranding — fondo transparente) ─────────────
# claro sin fondo.png = logo navy sobre transparente → slides CREMA
# oscuro sin fondo.png = logo blanco sobre transparente → slides DARK
LOGO_CREMA = r'C:\Users\MI PC\Desktop\PsicoEduca\Marketing\Rebranding\claro sin fondo.png'
LOGO_DARK_S= r'C:\Users\MI PC\Desktop\PsicoEduca\Marketing\Rebranding\oscuro sin fondo.png'
PHOTO      = r'C:\Users\MI PC\Desktop\PsicoEduca\Marketing\BOOK\sentado.JPG'

# ── PRESENTACION ────────────────────────────────────────
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]


# ── HELPERS ─────────────────────────────────────────────

def bg(slide, color=CREMA):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def txt(slide, text, x, y, w, h,
        font=SANS, size=16, bold=False, italic=False,
        color=DARK_T, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text = text
    r.font.name   = font
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def paras(slide, lines, x, y, w, h,
          font=SANS, size=14, color=DARK_T, align=PP_ALIGN.LEFT):
    """Multi-paragraph. line = str | dict{t,font,size,bold,italic,color}"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if isinstance(ln, dict):
            r = p.add_run()
            r.text = ln.get('t', '')
            r.font.name   = ln.get('font', font)
            r.font.size   = Pt(ln.get('size', size))
            r.font.bold   = ln.get('bold', False)
            r.font.italic = ln.get('italic', False)
            r.font.color.rgb = ln.get('color', color)
        else:
            r = p.add_run()
            r.text = str(ln)
            r.font.name  = font
            r.font.size  = Pt(size)
            r.font.color.rgb = color
    return tb


def logo(slide, dark_slide=False):
    """Logo esquina superior derecha, pequeño y discreto."""
    path = LOGO_DARK_S if dark_slide else LOGO_CREMA
    try:
        slide.shapes.add_picture(path, W - Cm(4.5), Cm(0.3), Cm(4.0), Cm(2.2))
    except Exception:
        pass  # si falla, continuar sin logo


def footer(slide, dark_slide=False):
    color = CREAM_T if dark_slide else RGBColor(0x88, 0x88, 0x88)
    txt(slide, 'Lic. Jean Clemotte  |  @Psico_Educa20',
        ML, H - Cm(1.1), Cm(22), Cm(0.8),
        SANS, 10, italic=True, color=color)


def box(slide, text, x, y, w, h,
        fill=RGBColor(0xE8, 0xF0, 0xFF), border=NAVY,
        font=SANS, size=14, bold=False, color=DARK_T,
        align=PP_ALIGN.LEFT):
    """Caja redondeada con texto."""
    s = slide.shapes.add_shape(5, x, y, w, h)   # 5 = rounded rect
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = border
    s.line.width = Pt(1)
    tf = s.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text = text
    r.font.name  = font
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.color.rgb = color
    return s


def big_number(slide, num_text, x, y, size=72, color=ORANGE):
    """Número o carácter grande como elemento visual."""
    txt(slide, num_text, x, y, Cm(5), Cm(5.5),
        SERIF, size, bold=True, color=color, align=PP_ALIGN.CENTER)


def accent_line(slide, x, y, w, color=ORANGE, thick=Pt(3)):
    """Línea de acento de color."""
    s = slide.shapes.add_shape(1, x, y, w, thick)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


def tbl(slide, data, x, y, w, h,
        hdr_bg=DARK_BG, hdr_fg=CREAM_T, odd=CREMA,
        even=RGBColor(0xE0, 0xDB, 0xD0), fs=12, center=False):
    rows = len(data)
    cols = max(len(r) for r in data)
    t = slide.shapes.add_table(rows, cols, x, y, w, h).table
    al = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
    for ri, row in enumerate(data):
        for ci in range(cols):
            val  = row[ci] if ci < len(row) else ''
            cell = t.cell(ri, ci)
            cell.text = str(val)
            tf = cell.text_frame
            tf.word_wrap = True
            for para in tf.paragraphs:
                para.alignment = al
                for run in para.runs:
                    run.font.name  = SANS
                    run.font.size  = Pt(fs)
                    run.font.bold  = (ri == 0)
                    run.font.color.rgb = hdr_fg if ri == 0 else DARK_T
            tcPr = cell._tc.get_or_add_tcPr()
            for old in tcPr.findall(qn('a:solidFill')):
                tcPr.remove(old)
            sf  = etree.SubElement(tcPr, qn('a:solidFill'))
            clr = etree.SubElement(sf,   qn('a:srgbClr'))
            if ri == 0:
                clr.set('val', '{:02X}{:02X}{:02X}'.format(*hdr_bg))
            elif ri % 2 == 1:
                clr.set('val', '{:02X}{:02X}{:02X}'.format(*odd))
            else:
                clr.set('val', '{:02X}{:02X}{:02X}'.format(*even))
    return t


def header_module(slide, module_label):
    """Etiqueta de módulo en caps en la parte superior, sobre fondo crema."""
    txt(slide, module_label.upper(),
        ML, Cm(0.5), Cm(22), Cm(0.7),
        SANS, 9, bold=True, color=NAVY)
    accent_line(slide, ML, Cm(1.3), Cm(22), color=NAVY, thick=Pt(1))


def bloque_slide(num_str, titulo, subtitulo=''):
    """Slide de transición de bloque — fondo dark."""
    sl = prs.slides.add_slide(blank)
    bg(sl, DARK_BG)
    big_number(sl, num_str, ML, Cm(3.5), size=80, color=ORANGE)
    txt(sl, titulo, ML + Cm(5.5), Cm(4.5), CW - Cm(6), Cm(4),
        SERIF, 40, bold=True, color=CREAM_T)
    if subtitulo:
        txt(sl, subtitulo, ML + Cm(5.5), Cm(9.0), CW - Cm(6), Cm(2),
            SANS, 18, color=SKY)
    logo(sl, dark_slide=True)
    footer(sl, dark_slide=True)
    return sl


# ════════════════════════════════════════════════════════
# S1 — PORTADA  (DARK)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl, DARK_BG)

# Acento visual: signo grande en naranja
txt(sl, '?', W - Cm(7), Cm(1.5), Cm(6), Cm(9),
    SERIF, 160, bold=True, color=RGBColor(0xE8, 0xA8, 0x35))
# Opacidad simulada con un rectángulo semi-transparente sobre el "?"
s = sl.shapes.add_shape(1, W - Cm(7), Cm(1.5), Cm(6), Cm(9))
s.fill.solid()
s.fill.fore_color.rgb = DARK_BG
s.line.fill.background()
# Set transparency via XML
from pptx.oxml.ns import qn as _qn; from lxml import etree as _et
xPr = s.fill._xPr; sf = xPr.solidFill
clr = sf.find(_qn('a:srgbClr'))
if clr is None:
    clr = _et.SubElement(sf, _qn('a:srgbClr')); clr.set('val','1E3A5F')
alpha = _et.SubElement(clr, _qn('a:alpha')); alpha.set('val','70000')

txt(sl, '¿Puede la mente\ncausar el\ncomportamiento?',
    ML, Cm(2.5), Cm(22), Cm(9),
    SERIF, 46, bold=True, color=CREAM_T)

accent_line(sl, ML, Cm(12.0), Cm(18), color=ORANGE, thick=Pt(3))

txt(sl, 'Bases filosóficas del Análisis de la Conducta',
    ML, Cm(12.8), Cm(22), Cm(1.5),
    SANS, 18, color=SKY)

txt(sl, 'Froxán et al. — Capítulo 3',
    ML, Cm(14.5), Cm(18), Cm(1.2),
    SANS, 14, italic=True, color=RGBColor(0xAA, 0xBB, 0xCC))

logo(sl, dark_slide=True)
footer(sl, dark_slide=True)


# ════════════════════════════════════════════════════════
# S2 — EL GANCHO: HOMERO  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl, CREMA)
header_module(sl, 'Bases filosóficas del análisis de la conducta')

txt(sl, 'Homero va al psicólogo',
    ML, Cm(1.7), CW - Cm(4.5), Cm(2.0),
    SERIF, 32, bold=True, color=DARK_T)

paras(sl, [
    {'t': '"El psicólogo dice: Homero tiene', 'size': 16, 'italic': True},
    {'t': 'personalidad impulsiva, baja autoestima', 'size': 16, 'italic': True},
    {'t': 'y conducta compulsiva."', 'size': 16, 'italic': True},
    ' ',
    {'t': '¿Eso explica por qué come donuts?', 'size': 17, 'bold': True},
    {'t': '¿Eso nos dice qué cambiar?', 'size': 17, 'bold': True},
], ML, Cm(4.2), Cm(18), Cm(9), SANS, 16, DARK_T)

# Elemento visual: "NO" grande en naranja
txt(sl, 'NO.', W - Cm(10), Cm(5.5), Cm(8), Cm(5),
    SERIF, 90, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
txt(sl, 'Solo le\npuso etiquetas.',
    W - Cm(10), Cm(10.5), Cm(8), Cm(3.5),
    SANS, 18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

accent_line(sl, ML, Cm(14.8), CW - Cm(1), color=ORANGE)
txt(sl, 'Esta pregunta va a guiar toda la presentacion.',
    ML, Cm(15.2), CW - Cm(1), Cm(1.2),
    SANS, 14, italic=True, color=NAVY)

logo(sl)
footer(sl)


# ════════════════════════════════════════════════════════
# S3 — ÍNDICE  (CREMA con elemento oscuro)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl, CREMA)
header_module(sl, 'Bases filosóficas del análisis de la conducta')

txt(sl, '¿Qué vamos a ver?',
    ML, Cm(1.7), Cm(18), Cm(2.0),
    SERIF, 34, bold=True, color=DARK_T)

bloques_idx = [
    ('01', 'El problema mente-cuerpo', 'Folk psychology  •  Error categorial'),
    ('02', '¿Qué significa "lo mental"?', 'Descriptivismo  •  Anti-descriptivismo  •  2 tipos de explicacion'),
    ('03', 'El AF como ciencia natural', 'Nivel agencial  •  Pensamiento como conducta'),
]
for i, (num, titulo, desc) in enumerate(bloques_idx):
    y = Cm(4.5) + i * Cm(4.0)
    # Número en dark con fondo
    s = sl.shapes.add_shape(5, ML, y, Cm(2.2), Cm(1.8))
    s.fill.solid(); s.fill.fore_color.rgb = DARK_BG; s.line.fill.background()
    tf = s.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num; r.font.name = SERIF
    r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = ORANGE
    # Título y descripción
    txt(sl, titulo, ML + Cm(2.8), y, Cm(25), Cm(1.0),
        SANS, 17, bold=True, color=DARK_T)
    txt(sl, desc, ML + Cm(2.8), y + Cm(1.1), Cm(25), Cm(0.9),
        SANS, 13, color=NAVY)
    if i < 2:
        accent_line(sl, ML, y + Cm(2.2), CW - Cm(1),
                    color=RGBColor(0xCC, 0xC8, 0xBF), thick=Pt(1))

txt(sl, 'Caso guia: Homero Simpson  —  Asuncion, Paraguay',
    ML, Cm(16.5), CW - Cm(1), Cm(1.0),
    SANS, 13, italic=True, color=NAVY)

logo(sl)
footer(sl)


# ════════════════════════════════════════════════════════
# S4 — PORTADA BLOQUE 1  (DARK)
# ════════════════════════════════════════════════════════
bloque_slide('01', 'El problema\nmente-cuerpo',
             'Folk psychology  •  Error categorial  •  Dos imágenes del mundo')


# ════════════════════════════════════════════════════════
# S5 — FOLK PSYCHOLOGY  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl, CREMA)
header_module(sl, 'Bloque 01 — El problema mente-cuerpo')

txt(sl, 'Así explicamos el comportamiento en la vida cotidiana',
    ML, Cm(1.7), CW - Cm(4.5), Cm(1.8),
    SERIF, 28, bold=True, color=DARK_T)

ejemplos = [
    ('"Homero come donuts', 'PORQUE le gustan y no tiene voluntad."'),
    ('"Lisa estudia ', 'PORQUE desea ser médica."'),
    ('"Bart molesta ', 'PORQUE es travieso y busca atención."'),
    ('"Marge aguanta ', 'PORQUE ama a su familia."'),
]
bw = CW / 2 - Cm(0.6)
for i, (parte1, parte2) in enumerate(ejemplos):
    col = i % 2
    row = i // 2
    x = ML + col * (bw + Cm(0.6))
    y = Cm(4.2) + row * Cm(3.2)
    s = sl.shapes.add_shape(5, x, y, bw, Cm(2.8))
    s.fill.solid(); s.fill.fore_color.rgb = RGBColor(0xE8, 0xF2, 0xFF)
    s.line.color.rgb = NAVY; s.line.width = Pt(1)
    tf = s.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run(); r1.text = parte1
    r1.font.name = SANS; r1.font.size = Pt(13); r1.font.color.rgb = DARK_T
    r2 = p.add_run(); r2.text = parte2
    r2.font.name = SANS; r2.font.size = Pt(13)
    r2.font.bold = True; r2.font.color.rgb = ORANGE

txt(sl, 'Esto es la "Folk Psychology" (Sellars, 1956):',
    ML, Cm(11.5), CW - Cm(1), Cm(0.9), SANS, 15, bold=True, color=NAVY)
txt(sl, 'explicar el comportamiento con CREENCIAS, DESEOS e INTENCIONES.',
    ML, Cm(12.5), CW - Cm(1), Cm(0.9), SANS, 15, color=DARK_T)
txt(sl, '¿Son estas explicaciones CIENTIFICAS? ¿Nos dicen las CAUSAS del comportamiento?',
    ML, Cm(13.8), CW - Cm(1), Cm(1.2), SANS, 14, italic=True, color=NAVY)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S6 — DOS IMÁGENES DEL MUNDO  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl, CREMA)
header_module(sl, 'Bloque 01 — El problema mente-cuerpo')

txt(sl, 'Dos formas de ver el mismo comportamiento',
    ML, Cm(1.7), CW - Cm(4.5), Cm(1.8),
    SERIF, 30, bold=True, color=DARK_T)

hw = CW / 2 - Cm(0.6)
# Columna izquierda: Imagen manifiesta
s = sl.shapes.add_shape(5, ML, Cm(4.0), hw, Cm(1.5))
s.fill.solid(); s.fill.fore_color.rgb = DARK_BG; s.line.fill.background()
tf = s.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = 'IMAGEN MANIFIESTA'
r.font.name = SANS; r.font.size = Pt(14); r.font.bold = True
r.font.color.rgb = CREAM_T

paras(sl, [
    {'t': '"Homero come donuts porque los desea', 'size': 14},
    {'t': ' y no tiene autocontrol."', 'size': 14},
    ' ',
    'Explica con RAZONES:',
    {'t': 'creencias, deseos, intenciones.', 'size': 14, 'color': ORANGE, 'bold': True},
    ' ',
    {'t': 'Util para evaluar si el comportamiento\nes racional o irracional.', 'size': 13, 'color': NAVY},
], ML, Cm(5.8), hw, Cm(8), SANS, 14, DARK_T)

# Columna derecha: Imagen científica
col2 = ML + hw + Cm(0.6)
s2 = sl.shapes.add_shape(5, col2, Cm(4.0), hw, Cm(1.5))
s2.fill.solid(); s2.fill.fore_color.rgb = NAVY; s2.line.fill.background()
tf2 = s2.text_frame; p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = 'IMAGEN CIENTIFICA'
r2.font.name = SANS; r2.font.size = Pt(14); r2.font.bold = True
r2.font.color.rgb = CREAM_T

paras(sl, [
    {'t': '"Homero come donuts porque comer donuts', 'size': 14},
    {'t': ' fue reforzado en su historia."', 'size': 14},
    ' ',
    'Explica con CAUSAS:',
    {'t': 'contingencias, historia de aprendizaje.', 'size': 14, 'color': GREEN_W, 'bold': True},
    ' ',
    {'t': 'Util para MODIFICAR el comportamiento.', 'size': 13, 'color': NAVY},
], col2, Cm(5.8), hw, Cm(8), SANS, 14, DARK_T)

accent_line(sl, ML, Cm(15.0), CW - Cm(1), ORANGE)
txt(sl, 'La ciencia del comportamiento necesita explicaciones CAUSALES, no solo evaluativas.',
    ML, Cm(15.4), CW - Cm(1), Cm(1.2), SANS, 14, bold=True, color=NAVY)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S7 — EL ERROR CATEGORIAL  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl, CREMA)
header_module(sl, 'Bloque 01 — El problema mente-cuerpo')

txt(sl, '"El fantasma en la máquina" — Gilbert Ryle (1949)',
    ML, Cm(1.7), CW - Cm(4.5), Cm(1.8),
    SERIF, 26, bold=True, color=DARK_T)

# Elemento visual: ícono de "X" grande indicando error
txt(sl, 'X', W - Cm(6.5), Cm(3.5), Cm(5), Cm(5),
    SERIF, 100, bold=True, color=RGBColor(0xDD, 0x44, 0x44),
    align=PP_ALIGN.CENTER)

data7 = [
    ['El error conceptual', 'Por que es un problema'],
    ['La mente se concibe como UNA COSA dentro del cuerpo\n(como si Homero tuviera un "motor de voluntad" adentro)',
     'La conducta no puede explicarse por algo que, por definición,\nnadie puede observar ni medir directamente'],
    ['"Homero tiene baja autoestima"\nsuena a describir una cosa real',
     'En realidad es una EVALUACION de su comportamiento:\ndecimos que su conducta se desvía de una norma'],
    ['La neurociencia busca "la voluntad" o "la autoestima"\nsubiendo a nivel subagencial',
     'Los conceptos mentales son irreducibles a estados cerebrales:\nson categorías distintas'],
]
tbl(sl, data7, ML, Cm(4.2), Cm(21), Cm(9.5), fs=13)

txt(sl, 'Concebir la mente como una "cosa" que CAUSA el comportamiento es un ERROR CATEGORIAL.',
    ML, Cm(14.5), CW - Cm(1), Cm(1.3), SANS, 15, bold=True, color=NAVY)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S8 — PORTADA BLOQUE 2  (DARK)
# ════════════════════════════════════════════════════════
bloque_slide('02', '¿Qué significa\n"lo mental"?',
             'Descriptivismo  •  Anti-descriptivismo  •  Wittgenstein')


# ════════════════════════════════════════════════════════
# S9 — DESCRIPTIVISMO Y SUS PROBLEMAS  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl, CREMA)
header_module(sl, 'Bloque 02 — ¿Que significa "lo mental"?')

txt(sl, 'El problema del descriptivismo',
    ML, Cm(1.7), CW - Cm(4.5), Cm(1.8),
    SERIF, 30, bold=True, color=DARK_T)

txt(sl, 'Si "Homero cree que los donuts son deliciosos"... ¿a qué HECHO se refiere esa oración?',
    ML, Cm(4.0), CW - Cm(1), Cm(1.2), SANS, 16, bold=True, color=NAVY)

hw9 = CW / 2 - Cm(0.6)
# Reduccionismo
s = sl.shapes.add_shape(5, ML, Cm(5.5), hw9, Cm(1.3))
s.fill.solid(); s.fill.fore_color.rgb = DARK_BG; s.line.fill.background()
tf = s.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
r = p.add_run(); r.text = 'REDUCCIONISMO'
r.font.name = SANS; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = ORANGE

paras(sl, [
    'Lo mental = estados cerebrales.',
    {'t': '"Baja serotonina → come donuts"', 'size': 14, 'italic': True},
    ' ',
    {'t': 'Problema:', 'bold': True, 'color': ORANGE},
    'Realización múltiple — distintos estados cerebrales\npueden producir la misma conducta.',
], ML, Cm(7.1), hw9, Cm(7), SANS, 14, DARK_T)

# Eliminativismo
col2 = ML + hw9 + Cm(0.6)
s2 = sl.shapes.add_shape(5, col2, Cm(5.5), hw9, Cm(1.3))
s2.fill.solid(); s2.fill.fore_color.rgb = NAVY; s2.line.fill.background()
tf2 = s2.text_frame; p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.LEFT
r2 = p2.add_run(); r2.text = 'ELIMINATIVISMO'
r2.font.name = SANS; r2.font.size = Pt(14); r2.font.bold = True
r2.font.color.rgb = GREEN_W

paras(sl, [
    'Lo mental no existe.',
    {'t': '"La autoestima es pseudociencia"', 'size': 14, 'italic': True},
    ' ',
    {'t': 'Problema:', 'bold': True, 'color': GREEN_W},
    'Pierde la capacidad de decir que Homero\nactua IRRACIONALMENTE.',
], col2, Cm(7.1), hw9, Cm(7), SANS, 14, DARK_T)

txt(sl, 'Ambas estrategias fracasan. Lo mental NO es un hecho que se puede describir como una silla o un átomo.',
    ML, Cm(15.0), CW - Cm(1), Cm(1.5), SANS, 14, italic=True, bold=True, color=NAVY)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S10 — ANTI-DESCRIPTIVISMO  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl, CREMA)
header_module(sl, 'Bloque 02 — ¿Que significa "lo mental"?')

txt(sl, 'Anti-descriptivismo: el significado depende del USO',
    ML, Cm(1.7), CW - Cm(4.5), Cm(1.8),
    SERIF, 28, bold=True, color=DARK_T)

# Cita grande de Wittgenstein
box(sl,
    '"El significado de cualquier expresión depende de las normas que rigen su uso,\nno de su capacidad para representar el mundo."\n\n— Wittgenstein, Investigaciones filosóficas',
    ML, Cm(4.0), CW - Cm(1), Cm(4.0),
    fill=DARK_BG, border=DARK_BG,
    font=SERIF, size=17, color=CREAM_T, align=PP_ALIGN.LEFT)
# Acento naranja en el borde izq de la cita
accent_line(sl, ML, Cm(4.0), Cm(0), color=ORANGE, thick=Pt(0))  # no funciona bien así, mejor rect
s_bar = sl.shapes.add_shape(1, ML, Cm(4.0), Pt(8), Cm(4.0))
s_bar.fill.solid(); s_bar.fill.fore_color.rgb = ORANGE; s_bar.line.fill.background()

txt(sl, 'Aplicado a Homero:',
    ML, Cm(8.8), CW - Cm(1), Cm(0.9), SANS, 16, bold=True, color=NAVY)

paras(sl, [
    {'t': '"Homero cree que los donuts son deliciosos"', 'size': 15, 'italic': True},
    {'t': 'no describe un objeto mental.', 'size': 15},
    {'t': 'Nos dice QUE ESPERAR de Homero:', 'size': 15, 'bold': True, 'color': ORANGE},
    {'t': '→ elegirá donuts, hablará de donuts, irá al Super Stock a buscarlos.', 'size': 14, 'italic': True},
    ' ',
    {'t': 'El significado vive en el COMPORTAMIENTO ESPERADO, no en su cabeza.', 'size': 15, 'bold': True, 'color': DARK_T},
], ML, Cm(9.8), CW - Cm(1), Cm(6), SANS, 15, DARK_T)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S11 — NORMATIVAS VS NOMOLÓGICAS  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl, CREMA)
header_module(sl, 'Bloque 02 — ¿Que significa "lo mental"?')

txt(sl, 'Dos tipos de explicación del comportamiento',
    ML, Cm(1.7), CW - Cm(4.5), Cm(1.8),
    SERIF, 30, bold=True, color=DARK_T)

hw11 = CW / 2 - Cm(0.6)
# NORMATIVA
s = sl.shapes.add_shape(5, ML, Cm(4.0), hw11, Cm(1.3))
s.fill.solid(); s.fill.fore_color.rgb = DARK_BG; s.line.fill.background()
tf = s.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = 'NORMATIVA  (razones)'
r.font.name = SANS; r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = ORANGE

paras(sl, [
    {'t': '"Homero no fue al trabajo porque CREYÓ que era feriado."', 'size': 14, 'italic': True},
    ' ',
    '→ Evalúa si el comportamiento es racional.',
    '→ NO explica la causa física.',
    {'t': 'Usa: creencias, deseos, intenciones.', 'bold': True, 'color': ORANGE, 'size': 14},
    {'t': 'Util para la ética y la vida cotidiana.', 'size': 13, 'italic': True, 'color': NAVY},
], ML, Cm(5.6), hw11, Cm(8.5), SANS, 14, DARK_T)

# NOMOLÓGICA
col2 = ML + hw11 + Cm(0.6)
s2 = sl.shapes.add_shape(5, col2, Cm(4.0), hw11, Cm(1.3))
s2.fill.solid(); s2.fill.fore_color.rgb = NAVY; s2.line.fill.background()
tf2 = s2.text_frame; p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = 'NOMOLÓGICA  (causas)'
r2.font.name = SANS; r2.font.size = Pt(15); r2.font.bold = True; r2.font.color.rgb = GREEN_W

paras(sl, [
    {'t': '"Homero no fue al trabajo porque faltar tuvo consecuencias reforzantes en el pasado."', 'size': 14, 'italic': True},
    ' ',
    '→ Establece las CAUSAS del comportamiento.',
    '→ Permite PREDECIR y MODIFICAR la conducta.',
    {'t': 'Usa: contingencias, historia de aprendizaje.', 'bold': True, 'color': GREEN_W, 'size': 14},
    {'t': 'Util para la intervención psicológica.', 'size': 13, 'italic': True, 'color': NAVY},
], col2, Cm(5.6), hw11, Cm(8.5), SANS, 14, DARK_T)

txt(sl, 'La psicología clínica necesita las DOS — pero sin confundirlas.',
    ML, Cm(15.0), CW - Cm(1), Cm(1.2), SANS, 15, bold=True, italic=True, color=NAVY)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S12 — PORTADA BLOQUE 3  (DARK)
# ════════════════════════════════════════════════════════
bloque_slide('03', 'El Análisis de la\nConducta como\nciencia natural',
             'Niveles de análisis  •  Pensamiento como conducta')


# ════════════════════════════════════════════════════════
# S13 — NIVEL AGENCIAL VS SUBAGENCIAL  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl, CREMA)
header_module(sl, 'Bloque 03 — El AF como ciencia natural')

txt(sl, '¿Desde qué nivel explicamos la conducta?',
    ML, Cm(1.7), CW - Cm(4.5), Cm(1.8),
    SERIF, 30, bold=True, color=DARK_T)

hw13 = CW / 2 - Cm(0.6)
s = sl.shapes.add_shape(5, ML, Cm(4.0), hw13, Cm(1.3))
s.fill.solid(); s.fill.fore_color.rgb = RGBColor(0x33, 0x55, 0x88)
s.line.fill.background()
tf = s.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = 'SUBAGENCIAL  (neurociencia)'
r.font.name = SANS; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = CREAM_T

paras(sl, [
    'Estudia lo que ocurre DENTRO del organismo.',
    {'t': 'Neuronas, dopamina, amígdala', 'color': SKY, 'bold': True},
    'Relaciones físico-contiguas.',
    ' ',
    {'t': 'Limitaciones:', 'bold': True, 'color': ORANGE},
    '• No permite intervención directa sobre neuronas',
    '• Realización múltiple: misma conducta, distintos estados cerebrales',
], ML, Cm(5.6), hw13, Cm(9), SANS, 14, DARK_T)

col2 = ML + hw13 + Cm(0.6)
s2 = sl.shapes.add_shape(5, col2, Cm(4.0), hw13, Cm(1.3))
s2.fill.solid(); s2.fill.fore_color.rgb = DARK_BG; s2.line.fill.background()
tf2 = s2.text_frame; p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = 'AGENCIAL  (análisis de conducta)'
r2.font.name = SANS; r2.font.size = Pt(14); r2.font.bold = True; r2.font.color.rgb = GREEN_W

paras(sl, [
    'Estudia la INTERACCION entre organismo y entorno.',
    {'t': 'Contingencias, historia de aprendizaje', 'color': GREEN_W, 'bold': True},
    'Relaciones temporales-funcionales.',
    ' ',
    {'t': 'Ventajas:', 'bold': True, 'color': GREEN_W},
    '• Intervención directa sobre contingencias',
    '• Predicción precisa del comportamiento',
], col2, Cm(5.6), hw13, Cm(9), SANS, 14, DARK_T)

txt(sl, 'Ambos son válidos y COMPLEMENTARIOS — pero no reducibles entre sí.',
    ML, Cm(15.3), CW - Cm(1), Cm(1.2), SANS, 14, bold=True, italic=True, color=NAVY)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S14 — EL PENSAMIENTO COMO CONDUCTA  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl, CREMA)
header_module(sl, 'Bloque 03 — El AF como ciencia natural')

txt(sl, 'El pensamiento no CAUSA la conducta — ES conducta',
    ML, Cm(1.7), CW - Cm(4.5), Cm(1.8),
    SERIF, 28, bold=True, color=DARK_T)

# Elemento visual: dos cajas conectadas con flecha
box(sl,
    'Visión mentalista:\n"Homero piensa: me merezco un donut"\n→ ese pensamiento CAUSA que coma',
    ML, Cm(4.2), Cm(13), Cm(4.0),
    fill=RGBColor(0xFF, 0xEE, 0xDD), border=RGBColor(0xDD, 0x66, 0x22),
    font=SANS, size=14, color=DARK_T)

# Flecha
txt(sl, '✗', ML + Cm(13.5), Cm(5.2), Cm(2), Cm(2),
    SANS, 40, bold=True, color=RGBColor(0xCC, 0x33, 0x33), align=PP_ALIGN.CENTER)

box(sl,
    'Conductismo radical:\nEl pensamiento encubierto ES conducta.\nNo causa la conducta — forma parte de ella.',
    ML, Cm(9.2), Cm(13), Cm(4.0),
    fill=RGBColor(0xE0, 0xFF, 0xF0), border=GREEN_W,
    font=SANS, size=14, color=DARK_T)

# Flecha
txt(sl, '✓', ML + Cm(13.5), Cm(10.2), Cm(2), Cm(2),
    SANS, 40, bold=True, color=GREEN_W, align=PP_ALIGN.CENTER)

paras(sl, [
    {'t': 'Implicación clínica:', 'bold': True, 'color': NAVY, 'size': 15},
    {'t': 'El autoinforme del paciente es un DATO, no la evidencia principal.', 'size': 14},
    {'t': 'Lo que Homero DICE puede no coincidir con lo que HACE.', 'size': 14, 'color': ORANGE},
], ML + Cm(15), Cm(5.5), Cm(14.5), Cm(8), SANS, 14, DARK_T)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S15 — VOLVEMOS A HOMERO  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl, CREMA)
header_module(sl, 'Bloque 03 — El AF como ciencia natural')

txt(sl, 'Homero: del diagnóstico al análisis filosófico-conductual',
    ML, Cm(1.7), CW - Cm(4.5), Cm(1.8),
    SERIF, 26, bold=True, color=DARK_T)

data15 = [
    ['Nivel de análisis', 'Explicación', 'Utilidad clínica'],
    ['Folk psychology\n(normativa)',
     '"Homero es impulsivo y tiene baja autoestima."',
     'Evalúa si es racional.\nNo permite intervenir.'],
    ['Neurociencia\n(subagencial)',
     '"Su sistema dopaminérgico responde al azúcar."',
     'Describe el sustrato biológico.\nNo indica cómo cambiar.'],
    ['Análisis de Conducta\n(agencial)',
     'Donut (Ed) → comer (R) → placer inmediato (C, RF+)\n+ evitar trabajo (RN)',
     'Explica el mantenimiento.\nDiseña la intervención.'],
    ['Conducta encubierta',
     '"Me lo merezco" = verbalización encubierta\nparte del patrón conductual',
     'También es objetivo de\nintervención.'],
]
tbl(sl, data15, ML, Cm(4.0), CW - Cm(1), Cm(11.5), fs=12)

txt(sl, 'El AF no dice que Homero "tiene" algo — explica POR QUÉ hace lo que hace y CÓMO cambiarlo.',
    ML, Cm(16.3), CW - Cm(1), Cm(1.3), SANS, 14, bold=True, color=NAVY)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S16 — IDEAS CLAVE  (DARK)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl, DARK_BG)

txt(sl, 'Para llevarse hoy',
    ML, Cm(1.5), CW - Cm(4.5), Cm(1.8),
    SERIF, 34, bold=True, color=CREAM_T, align=PP_ALIGN.CENTER)

ideas = [
    ('01', 'Las atribuciones mentales son NORMATIVAS — evalúan, no explican causas.\n"Tiene ansiedad" no es una explicación científica.'),
    ('02', 'El AF busca explicaciones NOMOLÓGICAS: contingencias que CAUSAN y MANTIENEN la conducta.'),
    ('03', 'El pensamiento no causa la conducta — es parte de ella.\nLa evidencia clínica es el COMPORTAMIENTO, no lo que el paciente dice.'),
]
for i, (num, texto) in enumerate(ideas):
    y = Cm(4.5) + i * Cm(4.0)
    # Número grande naranja
    s = sl.shapes.add_shape(5, ML, y, Cm(2.5), Cm(3.0))
    s.fill.solid(); s.fill.fore_color.rgb = ORANGE; s.line.fill.background()
    tf = s.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num
    r.font.name = SERIF; r.font.size = Pt(30); r.font.bold = True
    r.font.color.rgb = DARK_BG
    # Texto
    txt(sl, texto, ML + Cm(3.2), y + Cm(0.3), CW - Cm(7.5), Cm(2.8),
        SANS, 17, color=CREAM_T)

logo(sl, dark_slide=True)
footer(sl, dark_slide=True)


# ════════════════════════════════════════════════════════
# S17 — CIERRE  (DARK)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)

# Foto como fondo
try:
    sl.shapes.add_picture(PHOTO, 0, 0, W, H)
except Exception:
    pass
bg_shape = sl.shapes.add_shape(1, 0, 0, W*0.55, H)
bg_shape.fill.solid(); bg_shape.fill.fore_color.rgb = DARK_BG
bg_shape.line.fill.background()
xPr2 = bg_shape.fill._xPr; sf2 = xPr2.solidFill
c2 = sf2.find(_qn('a:srgbClr'))
if c2 is None:
    c2 = _et.SubElement(sf2, _qn('a:srgbClr')); c2.set('val','1E3A5F')
a2 = _et.SubElement(c2, _qn('a:alpha')); a2.set('val','85000')

accent_line(sl, ML, Cm(1.8), Cm(16), ORANGE, Pt(2))
accent_line(sl, ML, Cm(17.0), Cm(16), ORANGE, Pt(2))

txt(sl, '@PSICO_EDUCA20', ML, Cm(0.5), Cm(16), Cm(0.9),
    SANS, 12, color=RGBColor(0xAA, 0xBB, 0xCC))

txt(sl, 'Gracias por\nsu Atención',
    ML, Cm(3.5), Cm(16), Cm(7),
    SERIF, 52, bold=True, italic=True, color=CREAM_T)

txt(sl, 'Bases filosóficas del Análisis de la Conducta',
    ML, Cm(11.5), Cm(16), Cm(2), SANS, 17, color=SKY)

txt(sl, 'Lic. Jean Clemotte  |  PsicoEduca',
    ML, Cm(13.8), Cm(16), Cm(1.3), SANS, 15, italic=True, color=CREAM_T)

logo(sl, dark_slide=True)


# ════════════════════════════════════════════════════════
# GUARDAR
# ════════════════════════════════════════════════════════
output = r'C:\Users\MI PC\psicoeduca\materiales\presentacion-bases-filosoficas.pptx'
prs.save(output)
print('Listo: ' + output)
print('17 slides | nueva identidad PsicoEduca | Lic. Jean Clemotte')
