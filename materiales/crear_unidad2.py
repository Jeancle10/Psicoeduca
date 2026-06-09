"""
Presentacion Unidad 2: Cuestiones Filosóficas en torno al AC
Fuente: Cap. 3 — Frojan et al.
Marca: PsicoEduca — Rebranding
Autor: Lic. Jean Clemotte | @Psico_Educa20
Caso guia: Homero Simpson — Asuncion, Paraguay
"""
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
from pptx.oxml.ns import qn

# ── COLORES ─────────────────────────────────────────────
DARK_BG = RGBColor(0x1E, 0x3A, 0x5F)
CREMA   = RGBColor(0xF2, 0xED, 0xE4)
NAVY    = RGBColor(0x2B, 0x5E, 0xA7)
SKY     = RGBColor(0x4A, 0x9F, 0xE0)
ORANGE  = RGBColor(0xE8, 0xA8, 0x35)
GREEN_W = RGBColor(0x4A, 0xBF, 0xB0)
CREAM_T = RGBColor(0xF5, 0xF0, 0xDC)
DARK_T  = RGBColor(0x1E, 0x3A, 0x5F)
BOX_BG  = RGBColor(0xE4, 0xEE, 0xF8)
RED_C   = RGBColor(0xCC, 0x33, 0x33)
GREEN_C = RGBColor(0x22, 0x88, 0x44)

SERIF = 'Georgia'
SANS  = 'Calibri'
W  = Cm(33.87); H  = Cm(19.05)
ML = Cm(1.8);  CW = Cm(30.27)

LOGO_CREMA = r'C:\Users\MI PC\Desktop\PsicoEduca\Marketing\Rebranding\claro sin fondo.png'
LOGO_DARK  = r'C:\Users\MI PC\Desktop\PsicoEduca\Marketing\Rebranding\oscuro sin fondo.png'
PHOTO      = r'C:\Users\MI PC\Desktop\PsicoEduca\Marketing\BOOK\sentado.JPG'

prs = Presentation()
prs.slide_width = W; prs.slide_height = H
blank = prs.slide_layouts[6]

# ── HELPERS ─────────────────────────────────────────────
def bg(sl, c=CREMA):
    f = sl.background.fill; f.solid(); f.fore_color.rgb = c

def hbar(sl, x, y, w, color=ORANGE, t=Pt(3)):
    s = sl.shapes.add_shape(1, x, y, w, t)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

def txt(sl, text, x, y, w, h, font=SANS, size=16,
        bold=False, italic=False, color=DARK_T, align=PP_ALIGN.LEFT):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic; r.font.color.rgb = color
    return tb

def paras(sl, lines, x, y, w, h,
          font=SANS, size=14, color=DARK_T, align=PP_ALIGN.LEFT):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if isinstance(ln, dict):
            r = p.add_run(); r.text = ln.get('t', '')
            r.font.name   = ln.get('font', font)
            r.font.size   = Pt(ln.get('size', size))
            r.font.bold   = ln.get('bold', False)
            r.font.italic = ln.get('italic', False)
            r.font.color.rgb = ln.get('color', color)
        else:
            r = p.add_run(); r.text = str(ln)
            r.font.name = font; r.font.size = Pt(size); r.font.color.rgb = color

def logo(sl, dark=False):
    path = LOGO_DARK if dark else LOGO_CREMA
    try: sl.shapes.add_picture(path, W-Cm(4.5), Cm(0.3), Cm(4.0), Cm(2.2))
    except: pass

def footer(sl, dark=False):
    c = CREAM_T if dark else RGBColor(0x88, 0x88, 0x88)
    txt(sl, 'Lic. Jean Clemotte  |  @Psico_Educa20',
        ML, H-Cm(1.1), Cm(22), Cm(0.8), SANS, 10, italic=True, color=c)

def hmod(sl, label):
    txt(sl, label.upper(), ML, Cm(0.5), Cm(22), Cm(0.7), SANS, 9, bold=True, color=NAVY)
    hbar(sl, ML, Cm(1.3), Cm(22), NAVY, Pt(1))

def box(sl, text, x, y, w, h,
        fill=BOX_BG, border=NAVY, font=SANS, size=14,
        bold=False, color=DARK_T, align=PP_ALIGN.LEFT):
    s = sl.shapes.add_shape(5, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = border; s.line.width = Pt(1)
    tf = s.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = color
    return s

def tbl(sl, data, x, y, w, h,
        hdr_bg=DARK_BG, hdr_fg=CREAM_T,
        odd=CREMA, even=RGBColor(0xE0,0xDB,0xD0), fs=13):
    rows = len(data); cols = max(len(r) for r in data)
    t = sl.shapes.add_table(rows, cols, x, y, w, h).table
    for ri, row in enumerate(data):
        for ci in range(cols):
            val = row[ci] if ci < len(row) else ''
            cell = t.cell(ri, ci); cell.text = str(val)
            tf = cell.text_frame; tf.word_wrap = True
            for para in tf.paragraphs:
                para.alignment = PP_ALIGN.LEFT
                for run in para.runs:
                    run.font.name = SANS; run.font.size = Pt(fs)
                    run.font.bold = (ri == 0)
                    run.font.color.rgb = hdr_fg if ri == 0 else DARK_T
            tcPr = cell._tc.get_or_add_tcPr()
            for old in tcPr.findall(qn('a:solidFill')): tcPr.remove(old)
            sf = etree.SubElement(tcPr, qn('a:solidFill'))
            clr = etree.SubElement(sf, qn('a:srgbClr'))
            if ri == 0:   clr.set('val', '{:02X}{:02X}{:02X}'.format(*hdr_bg))
            elif ri%2==1: clr.set('val', '{:02X}{:02X}{:02X}'.format(*odd))
            else:         clr.set('val', '{:02X}{:02X}{:02X}'.format(*even))

def bloque(num, titulo, sub=''):
    sl = prs.slides.add_slide(blank); bg(sl, DARK_BG)
    txt(sl, num, ML, Cm(3.0), Cm(5), Cm(6),
        SERIF, 96, bold=True, color=ORANGE)
    hbar(sl, ML+Cm(5.8), Cm(6.5), Cm(20), ORANGE, Pt(3))
    txt(sl, titulo, ML+Cm(5.8), Cm(7.2), CW-Cm(7), Cm(5),
        SERIF, 38, bold=True, color=CREAM_T)
    if sub:
        txt(sl, sub, ML+Cm(5.8), Cm(12.5), CW-Cm(7), Cm(2),
            SANS, 16, color=SKY)
    logo(sl, dark=True); footer(sl, dark=True)


# ════════════════════════════════════════════════════════
# S1 — PORTADA  (DARK)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, DARK_BG)

txt(sl, '¿Puede la\nmente causar\nel comportamiento?',
    ML, Cm(1.8), Cm(21), Cm(10),
    SERIF, 46, bold=True, color=CREAM_T)

hbar(sl, ML, Cm(12.2), Cm(20), ORANGE, Pt(3))
txt(sl, 'Unidad 2 — Bases filosóficas del Análisis de la Conducta',
    ML, Cm(12.8), Cm(25), Cm(1.4), SANS, 16, color=SKY)
txt(sl, 'Froján et al. — Capítulo 3',
    ML, Cm(14.4), Cm(18), Cm(1.1),
    SERIF, 14, italic=True, color=RGBColor(0xAA,0xBB,0xCC))

# Visual: gran "?" en el lado derecho
txt(sl, '?', W-Cm(10), Cm(0.5), Cm(9), Cm(17),
    SERIF, 200, bold=True, color=RGBColor(0x2B,0x4A,0x6E),
    align=PP_ALIGN.CENTER)

logo(sl, dark=True); footer(sl, dark=True)


# ════════════════════════════════════════════════════════
# S2 — GANCHO: HOMERO  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, CREMA)
hmod(sl, 'Unidad 2 — Bases filosóficas del AC')

txt(sl, 'Homero va al psicólogo (porque Marge lo mandó)',
    ML, Cm(1.7), CW-Cm(5), Cm(2.0),
    SERIF, 28, bold=True, color=DARK_T)

paras(sl, [
    {'t': '"El psicólogo diagnostica: personalidad impulsiva,', 'size': 16, 'italic': True},
    {'t': ' baja autoestima, falta de voluntad."', 'size': 16, 'italic': True},
    ' ',
    {'t': '¿Eso explica POR QUÉ Homero come donuts?', 'size': 17, 'bold': True, 'color': NAVY},
    {'t': '¿Eso nos dice QUÉ hay que cambiar?', 'size': 17, 'bold': True, 'color': NAVY},
], ML, Cm(4.3), Cm(19), Cm(8), SANS, 16, DARK_T)

# Visual: "NO." naranja grande
txt(sl, 'NO.',
    Cm(21.5), Cm(4.5), Cm(10.5), Cm(5.5),
    SERIF, 90, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
txt(sl, 'Solo le\npuso etiquetas.',
    Cm(21.5), Cm(10.2), Cm(10.5), Cm(3),
    SANS, 19, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

hbar(sl, ML, Cm(14.0), CW-Cm(1), ORANGE)
txt(sl, 'Esta unidad explica POR QUÉ esas explicaciones no son científicas y qué propone el AF en su lugar.',
    ML, Cm(14.4), CW-Cm(1), Cm(1.3),
    SANS, 14, italic=True, color=NAVY)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S3 — ÍNDICE  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, CREMA)
hmod(sl, 'Unidad 2 — Bases filosóficas del AC')

txt(sl, '¿Qué vamos a ver?',
    ML, Cm(1.7), Cm(12), Cm(2.0),
    SERIF, 34, bold=True, color=DARK_T)

items = [
    ('01', 'El problema mente-cuerpo',
     'Folk psychology  •  Error categorial  •  Descriptivismo'),
    ('02', '¿Qué significa "lo mental"?',
     'Anti-descriptivismo (Wittgenstein)  •  Normativas vs Nomológicas'),
    ('03', 'El AF como ciencia natural',
     'Niveles de análisis  •  Nexos temporales  •  Pensamiento como conducta'),
]
for i, (num, tit, desc) in enumerate(items):
    y = Cm(4.5) + i * Cm(4.0)
    s = sl.shapes.add_shape(5, ML, y, Cm(2.3), Cm(1.9))
    s.fill.solid(); s.fill.fore_color.rgb = DARK_BG; s.line.fill.background()
    tf = s.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num
    r.font.name = SERIF; r.font.size = Pt(22); r.font.bold = True
    r.font.color.rgb = ORANGE
    txt(sl, tit, ML+Cm(3.0), y, CW-Cm(4), Cm(1.1),
        SANS, 17, bold=True, color=DARK_T)
    txt(sl, desc, ML+Cm(3.0), y+Cm(1.2), CW-Cm(4), Cm(0.9),
        SANS, 13, color=NAVY)
    if i < 2:
        hbar(sl, ML, y+Cm(2.3), CW-Cm(1), RGBColor(0xCC,0xC8,0xBF), Pt(1))

txt(sl, 'Caso guía: Homero Simpson  —  Asunción, Paraguay',
    ML, Cm(16.5), CW-Cm(1), Cm(1.0),
    SANS, 13, italic=True, color=NAVY)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S4 — PORTADA BLOQUE 1  (DARK)
# ════════════════════════════════════════════════════════
bloque('01', 'El problema\nmente-cuerpo',
       'Folk psychology  •  Imagen manifiesta vs científica  •  Error categorial')


# ════════════════════════════════════════════════════════
# S5 — FOLK PSYCHOLOGY  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, CREMA)
hmod(sl, 'Bloque 01 — El problema mente-cuerpo')

txt(sl, 'Así explicamos el comportamiento en la vida cotidiana',
    ML, Cm(1.7), CW-Cm(5), Cm(1.8),
    SERIF, 28, bold=True, color=DARK_T)

ejemplos = [
    ('"Homero come donuts ', 'PORQUE le gustan y no tiene voluntad."'),
    ('"Lisa estudia ', 'PORQUE desea ser médica."'),
    ('"Bart molesta ', 'PORQUE es travieso y busca atención."'),
    ('"Marge aguanta ', 'PORQUE ama a su familia."'),
]
bw = CW/2 - Cm(0.6)
for i, (p1, p2) in enumerate(ejemplos):
    col = i%2; row = i//2
    x = ML + col*(bw+Cm(0.6)); y = Cm(4.2) + row*Cm(3.2)
    s = sl.shapes.add_shape(5, x, y, bw, Cm(2.8))
    s.fill.solid(); s.fill.fore_color.rgb = BOX_BG
    s.line.color.rgb = NAVY; s.line.width = Pt(1)
    tf = s.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run(); r1.text = p1
    r1.font.name = SANS; r1.font.size = Pt(14); r1.font.color.rgb = DARK_T
    r2 = p.add_run(); r2.text = p2
    r2.font.name = SANS; r2.font.size = Pt(14)
    r2.font.bold = True; r2.font.color.rgb = ORANGE

paras(sl, [
    {'t': 'Esto es la "Folk Psychology" (Sellars, 1956):', 'bold': True, 'size': 15, 'color': NAVY},
    'explicar el comportamiento con CREENCIAS, DESEOS e INTENCIONES.',
    ' ',
    {'t': '¿Son explicaciones científicas? ¿Nos dicen las CAUSAS?', 'size': 14, 'italic': True, 'color': NAVY},
], ML, Cm(11.5), CW-Cm(1), Cm(4), SANS, 15, DARK_T)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S6 — DOS IMÁGENES DEL MUNDO  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, CREMA)
hmod(sl, 'Bloque 01 — El problema mente-cuerpo')

txt(sl, 'Dos formas de ver el mismo comportamiento (Sellars, 1956)',
    ML, Cm(1.7), CW-Cm(5), Cm(1.8),
    SERIF, 28, bold=True, color=DARK_T)

hw = CW/2 - Cm(0.6)

# Imagen manifiesta
s1 = sl.shapes.add_shape(5, ML, Cm(4.0), hw, Cm(1.4))
s1.fill.solid(); s1.fill.fore_color.rgb = DARK_BG; s1.line.fill.background()
tf = s1.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = 'IMAGEN MANIFIESTA  (folk psychology)'
r.font.name = SANS; r.font.size = Pt(13); r.font.bold = True
r.font.color.rgb = ORANGE

paras(sl, [
    {'t': '"Homero come donuts porque los desea\n y no tiene autocontrol."', 'italic': True, 'size': 15},
    ' ',
    'Explica con RAZONES:',
    {'t': 'creencias, deseos, intenciones.', 'bold': True, 'color': ORANGE},
    ' ',
    {'t': 'Útil para evaluar si el comportamiento\nes racional o irracional.', 'size': 13, 'italic': True, 'color': NAVY},
], ML, Cm(5.7), hw, Cm(9), SANS, 14, DARK_T)

# Imagen científica
col2 = ML + hw + Cm(0.6)
s2 = sl.shapes.add_shape(5, col2, Cm(4.0), hw, Cm(1.4))
s2.fill.solid(); s2.fill.fore_color.rgb = NAVY; s2.line.fill.background()
tf2 = s2.text_frame; p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = 'IMAGEN CIENTÍFICA  (análisis de conducta)'
r2.font.name = SANS; r2.font.size = Pt(13); r2.font.bold = True
r2.font.color.rgb = GREEN_W

paras(sl, [
    {'t': '"Homero come donuts porque comer donuts fue\n reforzado en su historia."', 'italic': True, 'size': 15},
    ' ',
    'Explica con CAUSAS:',
    {'t': 'contingencias, historia de aprendizaje.', 'bold': True, 'color': GREEN_W},
    ' ',
    {'t': 'Útil para MODIFICAR el comportamiento.', 'size': 13, 'italic': True, 'color': NAVY},
], col2, Cm(5.7), hw, Cm(9), SANS, 14, DARK_T)

hbar(sl, ML, Cm(15.3), CW-Cm(1), ORANGE)
txt(sl, 'La psicología como ciencia necesita explicaciones CAUSALES, no solo evaluativas.',
    ML, Cm(15.7), CW-Cm(1), Cm(1.2),
    SANS, 15, bold=True, color=NAVY)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S7 — ERROR CATEGORIAL  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, CREMA)
hmod(sl, 'Bloque 01 — El problema mente-cuerpo')

txt(sl, '"El fantasma en la máquina" — Gilbert Ryle (1949)',
    ML, Cm(1.7), CW-Cm(5), Cm(1.8),
    SERIF, 26, bold=True, color=DARK_T)

# Visual: X grande de error
txt(sl, '✗', W-Cm(7.5), Cm(3.0), Cm(6), Cm(6),
    SANS, 110, bold=True, color=RED_C, align=PP_ALIGN.CENTER)

paras(sl, [
    {'t': 'El error de Descartes:', 'bold': True, 'size': 17, 'color': NAVY},
    'Concebir la mente como UNA COSA dentro del cuerpo.',
    '→ "Homero tiene un motor de voluntad adentro"',
    ' ',
    {'t': 'El problema:', 'bold': True, 'size': 15, 'color': ORANGE},
    'La conducta no puede explicarse por algo que',
    {'t': 'nadie puede observar ni medir directamente.', 'italic': True},
    ' ',
    {'t': '"Homero tiene baja autoestima" no describe', 'size': 14},
    {'t': 'una cosa real — es una EVALUACIÓN de su conducta.', 'size': 14, 'bold': True, 'color': ORANGE},
], ML, Cm(4.2), Cm(22), Cm(11), SANS, 15, DARK_T)

hbar(sl, ML, Cm(15.8), CW-Cm(1), GREEN_W)
txt(sl, 'Concebir la mente como una "cosa" que CAUSA el comportamiento es un ERROR CATEGORIAL.',
    ML, Cm(16.2), CW-Cm(1), Cm(1.3),
    SANS, 15, bold=True, color=NAVY)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S8 — PORTADA BLOQUE 2  (DARK)
# ════════════════════════════════════════════════════════
bloque('02', '¿Qué significa\n"lo mental"?',
       'Descriptivismo  •  Anti-descriptivismo  •  2 tipos de explicación')


# ════════════════════════════════════════════════════════
# S9 — DESCRIPTIVISMO: REDUCCIONISMO Y ELIMINATIVISMO  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, CREMA)
hmod(sl, 'Bloque 02 — ¿Qué significa "lo mental"?')

txt(sl, 'El problema del descriptivismo',
    ML, Cm(1.7), CW-Cm(5), Cm(1.8),
    SERIF, 30, bold=True, color=DARK_T)

txt(sl, '"Homero cree que los donuts son deliciosos"  →  ¿a qué HECHO se refiere esa oración?',
    ML, Cm(4.0), CW-Cm(1), Cm(1.1),
    SANS, 16, bold=True, color=NAVY)

hw9 = CW/2 - Cm(0.6)

s1 = sl.shapes.add_shape(5, ML, Cm(5.5), hw9, Cm(1.3))
s1.fill.solid(); s1.fill.fore_color.rgb = DARK_BG; s1.line.fill.background()
tf = s1.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
r = p.add_run(); r.text = 'REDUCCIONISMO'
r.font.name = SANS; r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = ORANGE

paras(sl, [
    'Lo mental = estados cerebrales.',
    {'t': '"Baja serotonina → come donuts"', 'italic': True, 'size': 14},
    ' ',
    {'t': 'Problema:', 'bold': True, 'color': ORANGE},
    'Realización múltiple: distintos estados cerebrales',
    'pueden producir la misma conducta.',
], ML, Cm(7.1), hw9, Cm(8.5), SANS, 14, DARK_T)

col2 = ML + hw9 + Cm(0.6)
s2 = sl.shapes.add_shape(5, col2, Cm(5.5), hw9, Cm(1.3))
s2.fill.solid(); s2.fill.fore_color.rgb = NAVY; s2.line.fill.background()
tf2 = s2.text_frame; p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.LEFT
r2 = p2.add_run(); r2.text = 'ELIMINATIVISMO'
r2.font.name = SANS; r2.font.size = Pt(15); r2.font.bold = True; r2.font.color.rgb = GREEN_W

paras(sl, [
    'Lo mental no existe.',
    {'t': '"La autoestima es pseudociencia"', 'italic': True, 'size': 14},
    ' ',
    {'t': 'Problema:', 'bold': True, 'color': GREEN_W},
    'Perdemos la capacidad de decir que Homero',
    'actúa IRRACIONALMENTE.',
], col2, Cm(7.1), hw9, Cm(8.5), SANS, 14, DARK_T)

hbar(sl, ML, Cm(16.2), CW-Cm(1), ORANGE)
txt(sl, 'Ambas fracasan. Lo mental NO es un hecho que se puede describir como una silla o un átomo.',
    ML, Cm(16.6), CW-Cm(1), Cm(1.2),
    SANS, 14, italic=True, bold=True, color=NAVY)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S10 — ANTI-DESCRIPTIVISMO: WITTGENSTEIN  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, CREMA)
hmod(sl, 'Bloque 02 — ¿Qué significa "lo mental"?')

txt(sl, 'Anti-descriptivismo: el significado depende del USO',
    ML, Cm(1.7), CW-Cm(5), Cm(1.8),
    SERIF, 28, bold=True, color=DARK_T)

# Cita en caja dark con barra naranja lateral
s_bar = sl.shapes.add_shape(1, ML, Cm(4.0), Pt(8), Cm(4.0))
s_bar.fill.solid(); s_bar.fill.fore_color.rgb = ORANGE; s_bar.line.fill.background()

box(sl,
    '"El significado de cualquier expresión depende de las normas\n'
    'que rigen su uso, no de su capacidad para representar el mundo."\n\n'
    '— Wittgenstein, Investigaciones filosóficas (1953)',
    ML + Cm(0.7), Cm(4.0), CW-Cm(1.5), Cm(4.0),
    fill=DARK_BG, border=DARK_BG, font=SERIF, size=16, color=CREAM_T)

txt(sl, 'Aplicado a Homero:',
    ML, Cm(9.0), CW-Cm(1), Cm(0.9), SANS, 16, bold=True, color=NAVY)

paras(sl, [
    {'t': '"Homero cree que los donuts son deliciosos"', 'size': 16, 'italic': True},
    {'t': 'no describe un objeto mental en su cabeza.', 'size': 15},
    ' ',
    {'t': 'Nos dice QUÉ ESPERAR de Homero:', 'bold': True, 'size': 15, 'color': ORANGE},
    '→ elegirá donuts, hablará de donuts, irá al Super Stock a buscarlos.',
    ' ',
    {'t': 'El significado vive en el COMPORTAMIENTO ESPERADO.', 'bold': True, 'size': 15},
], ML, Cm(10.0), CW-Cm(1), Cm(6.5), SANS, 15, DARK_T)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S11 — NORMATIVAS VS NOMOLÓGICAS  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, CREMA)
hmod(sl, 'Bloque 02 — ¿Qué significa "lo mental"?')

txt(sl, 'Dos tipos de explicación del comportamiento',
    ML, Cm(1.7), CW-Cm(5), Cm(1.8),
    SERIF, 30, bold=True, color=DARK_T)

hw11 = CW/2 - Cm(0.6)

s1 = sl.shapes.add_shape(5, ML, Cm(4.0), hw11, Cm(1.4))
s1.fill.solid(); s1.fill.fore_color.rgb = DARK_BG; s1.line.fill.background()
tf = s1.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = 'NORMATIVA  (razones)'
r.font.name = SANS; r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = ORANGE

paras(sl, [
    {'t': '"Homero no fue al trabajo porque\n CREYÓ que era feriado."', 'italic': True, 'size': 15},
    ' ',
    '→ Evalúa si el comportamiento es racional.',
    {'t': 'Usa: creencias, deseos, intenciones.', 'bold': True, 'color': ORANGE},
    {'t': 'Útil para la ética y la vida cotidiana.', 'size': 13, 'italic': True, 'color': NAVY},
], ML, Cm(5.7), hw11, Cm(9.5), SANS, 14, DARK_T)

col2 = ML + hw11 + Cm(0.6)
s2 = sl.shapes.add_shape(5, col2, Cm(4.0), hw11, Cm(1.4))
s2.fill.solid(); s2.fill.fore_color.rgb = NAVY; s2.line.fill.background()
tf2 = s2.text_frame; p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = 'NOMOLÓGICA  (causas)'
r2.font.name = SANS; r2.font.size = Pt(15); r2.font.bold = True; r2.font.color.rgb = GREEN_W

paras(sl, [
    {'t': '"Homero no fue al trabajo porque faltar\n tuvo consecuencias reforzantes."', 'italic': True, 'size': 15},
    ' ',
    '→ Establece las CAUSAS del comportamiento.',
    {'t': 'Usa: contingencias, historia de aprendizaje.', 'bold': True, 'color': GREEN_W},
    {'t': 'Útil para la intervención psicológica.', 'size': 13, 'italic': True, 'color': NAVY},
], col2, Cm(5.7), hw11, Cm(9.5), SANS, 14, DARK_T)

hbar(sl, ML, Cm(16.0), CW-Cm(1), GREEN_W)
txt(sl, 'La psicología clínica necesita las DOS — pero sin confundirlas.',
    ML, Cm(16.4), CW-Cm(1), Cm(1.2),
    SANS, 15, bold=True, italic=True, color=NAVY)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S12 — PORTADA BLOQUE 3  (DARK)
# ════════════════════════════════════════════════════════
bloque('03', 'El AF como\nciencia natural',
       'Niveles de análisis  •  Nexos temporales  •  Pensamiento como conducta')


# ════════════════════════════════════════════════════════
# S13 — NIVEL AGENCIAL VS SUBAGENCIAL  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, CREMA)
hmod(sl, 'Bloque 03 — El AF como ciencia natural')

txt(sl, '¿Desde qué nivel explicamos la conducta?',
    ML, Cm(1.7), CW-Cm(5), Cm(1.8),
    SERIF, 30, bold=True, color=DARK_T)

hw13 = CW/2 - Cm(0.6)

s1 = sl.shapes.add_shape(5, ML, Cm(4.2), hw13, Cm(1.4))
s1.fill.solid(); s1.fill.fore_color.rgb = RGBColor(0x33,0x55,0x88)
s1.line.fill.background()
tf = s1.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = 'SUBAGENCIAL  (neurociencia)'
r.font.name = SANS; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = CREAM_T

paras(sl, [
    'Estudia lo que ocurre DENTRO del organismo.',
    {'t': 'Neuronas, dopamina, amígdala', 'color': SKY, 'bold': True},
    'Relaciones de contigüidad FÍSICA.',
    ' ',
    {'t': 'Limitaciones:', 'bold': True, 'color': ORANGE},
    '• No permite intervención directa',
    '• Realización múltiple: misma conducta,\n  distintos estados cerebrales',
], ML, Cm(5.9), hw13, Cm(9.5), SANS, 14, DARK_T)

col2 = ML + hw13 + Cm(0.6)
s2 = sl.shapes.add_shape(5, col2, Cm(4.2), hw13, Cm(1.4))
s2.fill.solid(); s2.fill.fore_color.rgb = DARK_BG; s2.line.fill.background()
tf2 = s2.text_frame; p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = 'AGENCIAL  (análisis de conducta)'
r2.font.name = SANS; r2.font.size = Pt(14); r2.font.bold = True; r2.font.color.rgb = GREEN_W

paras(sl, [
    'Estudia la INTERACCIÓN organismo-entorno.',
    {'t': 'Contingencias, historia de aprendizaje', 'color': GREEN_W, 'bold': True},
    'Relaciones de contigüidad TEMPORAL.',
    ' ',
    {'t': 'Ventajas:', 'bold': True, 'color': GREEN_W},
    '• Intervención directa sobre contingencias',
    '• Predicción precisa del comportamiento',
], col2, Cm(5.9), hw13, Cm(9.5), SANS, 14, DARK_T)

txt(sl, 'Ambos son válidos y COMPLEMENTARIOS — pero no reducibles entre sí.',
    ML, Cm(16.2), CW-Cm(1), Cm(1.2),
    SANS, 14, bold=True, italic=True, color=NAVY)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S14 — EL AF COMO DARWIN: NEXOS TEMPORALES  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, CREMA)
hmod(sl, 'Bloque 03 — El AF como ciencia natural')

txt(sl, 'El AF es la "selección natural" del comportamiento individual',
    ML, Cm(1.7), CW-Cm(5), Cm(1.8),
    SERIF, 26, bold=True, color=DARK_T)

hw14 = CW/2 - Cm(0.7)

box(sl,
    'SELECCIÓN NATURAL\n(evolución de la especie)\n\n'
    'El ambiente selecciona rasgos por sus CONSECUENCIAS\na lo largo de la historia de la especie.\n\n'
    'Nexo causal: TEMPORAL — no mecánico.',
    ML, Cm(4.2), hw14, Cm(5.5),
    fill=BOX_BG, border=NAVY, size=15)

# Flecha entre columnas
txt(sl, '↔', ML + hw14 + Cm(0.2), Cm(5.8), Cm(1.3), Cm(2.5),
    SANS, 36, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

box(sl,
    'SELECCIÓN POR CONTINGENCIAS\n(historia individual)\n\n'
    'Las conductas con consecuencias reforzantes\nSE REPITEN a lo largo de la vida de Homero.\n\n'
    'Nexo causal: TEMPORAL — no mecánico.',
    ML + hw14 + Cm(1.5), Cm(4.2), hw14, Cm(5.5),
    fill=BOX_BG, border=DARK_BG, size=15)

paras(sl, [
    {'t': 'El AF no busca nexos físico-contiguos (neurona → conducta)', 'size': 15, 'bold': True},
    {'t': 'sino nexos TEMPORALES-FUNCIONALES (contingencia → conducta).', 'size': 15, 'bold': True, 'color': ORANGE},
    ' ',
    {'t': 'Homero come donuts no se explica por sus neuronas,', 'size': 14},
    {'t': 'sino por las CONSECUENCIAS de comer donuts en su historia.', 'size': 14, 'italic': True},
], ML, Cm(10.5), CW-Cm(1), Cm(5.5), SANS, 15, DARK_T)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S15 — EL PENSAMIENTO COMO CONDUCTA  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, CREMA)
hmod(sl, 'Bloque 03 — El AF como ciencia natural')

txt(sl, 'El pensamiento no CAUSA la conducta — ES conducta',
    ML, Cm(1.7), CW-Cm(5), Cm(1.8),
    SERIF, 28, bold=True, color=DARK_T)

data15 = [
    ['', 'Visión mentalista', 'Conductismo radical'],
    ['¿Qué es el\npensamiento?',
     'Un proceso INTERNO que CAUSA el comportamiento.\n"Homero piensa que merece un donut → come."',
     'Conducta verbal ENCUBIERTA.\nParte del patrón conductual, no su causa.'],
    ['¿Qué son las\nemociones?',
     'Estados mentales que causan la conducta.\n"Homero come por ansiedad."',
     'También son conducta encubierta.\nSe explican por contingencias, no explican la conducta.'],
    ['¿Qué dice\nel psicólogo?',
     '"Cambiemos sus pensamientos y cambiará su conducta."\n(terapia cognitiva clásica)',
     '"El autoinforme es un dato, no la evidencia principal.\nLo que Homero HACE es lo que importa."'],
]
tbl(sl, data15, ML, Cm(4.2), CW-Cm(1), Cm(10.5), fs=13)

hbar(sl, ML, Cm(15.5), CW-Cm(1), ORANGE)
txt(sl, 'Implicación clínica: lo que Homero DICE puede no coincidir con lo que HACE. El comportamiento es la evidencia.',
    ML, Cm(15.9), CW-Cm(1), Cm(1.3),
    SANS, 14, bold=True, color=NAVY)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S16 — VOLVEMOS A HOMERO  (CREMA)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, CREMA)
hmod(sl, 'Bloque 03 — El AF como ciencia natural')

txt(sl, 'Homero: del diagnóstico al análisis filosófico-conductual',
    ML, Cm(1.7), CW-Cm(5), Cm(1.8),
    SERIF, 24, bold=True, color=DARK_T)

data16 = [
    ['Nivel', 'Explicación de Homero', 'Tipo de explicación'],
    ['Folk psychology',
     '"Homero es impulsivo y tiene baja autoestima."',
     'NORMATIVA: evalúa\nNo permite intervenir'],
    ['Neurociencia',
     '"Su sistema de recompensa responde al azúcar."',
     'NOMOLÓGICA subagencial:\ndescribe sustrato, no interviene'],
    ['Análisis de Conducta',
     'Donut disponible (Ed) → comer (R) → placer inmediato (RF+)\n+ evitar pensamientos de dieta (RN)',
     'NOMOLÓGICA agencial:\nexplica, predice, interviene'],
    ['Conducta encubierta',
     '"Me lo merezco" / "total uno más no cambia nada" =\nverbalización encubierta, parte del patrón',
     'También es objetivo de\nintervención, no causa'],
]
tbl(sl, data16, ML, Cm(4.0), CW-Cm(1), Cm(11.0), fs=13)

hbar(sl, ML, Cm(15.8), CW-Cm(1), GREEN_W)
txt(sl, 'El AF no dice que Homero "tiene" algo — explica POR QUÉ hace lo que hace y CÓMO cambiarlo.',
    ML, Cm(16.2), CW-Cm(1), Cm(1.2),
    SANS, 14, bold=True, color=NAVY)

logo(sl); footer(sl)


# ════════════════════════════════════════════════════════
# S17 — IDEAS CLAVE  (DARK)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, DARK_BG)

txt(sl, 'Para llevarse hoy',
    ML, Cm(1.5), CW-Cm(5), Cm(1.8),
    SERIF, 36, bold=True, color=CREAM_T)

ideas = [
    ('01', ORANGE,  'Las atribuciones mentales son NORMATIVAS — evalúan, no explican causas.\n"Tiene ansiedad" o "es impulsivo" no son explicaciones científicas.'),
    ('02', GREEN_W, 'El AF busca explicaciones NOMOLÓGICAS en el nivel AGENCIAL:\nlas contingencias que CAUSAN y MANTIENEN la conducta.'),
    ('03', ORANGE,  'El pensamiento no causa la conducta — es parte de ella.\nLa evidencia clínica es el COMPORTAMIENTO, no el autoinforme.'),
]
for i, (num, accent, texto) in enumerate(ideas):
    y = Cm(4.5) + i * Cm(4.0)
    s = sl.shapes.add_shape(5, ML, y, Cm(2.5), Cm(3.0))
    s.fill.solid(); s.fill.fore_color.rgb = accent; s.line.fill.background()
    tf = s.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num
    r.font.name = SERIF; r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = DARK_BG
    txt(sl, texto, ML+Cm(3.2), y+Cm(0.3), CW-Cm(7.5), Cm(2.8),
        SANS, 17, color=CREAM_T)

logo(sl, dark=True); footer(sl, dark=True)


# ════════════════════════════════════════════════════════
# S18 — CIERRE CON FOTO  (DARK)
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
try:
    sl.shapes.add_picture(PHOTO, 0, 0, W, H)
except: pass

overlay = sl.shapes.add_shape(1, 0, 0, W*0.54, H)
overlay.fill.solid(); overlay.fill.fore_color.rgb = DARK_BG
overlay.line.fill.background()
from pptx.oxml.ns import qn as _q; from lxml import etree as _e
xPr = overlay.fill._xPr; sf = xPr.solidFill
c = sf.find(_q('a:srgbClr'))
if c is None:
    c = _e.SubElement(sf, _q('a:srgbClr')); c.set('val','1E3A5F')
a = _e.SubElement(c, _q('a:alpha')); a.set('val','85000')

hbar(sl, ML, Cm(1.8), Cm(16), ORANGE, Pt(2))
hbar(sl, ML, Cm(17.0), Cm(16), ORANGE, Pt(2))
txt(sl, '@PSICO_EDUCA20', ML, Cm(0.5), Cm(16), Cm(0.9),
    SANS, 12, color=RGBColor(0xAA,0xBB,0xCC))
txt(sl, 'Gracias por\nsu Atención',
    ML, Cm(3.5), Cm(16), Cm(7),
    SERIF, 52, bold=True, italic=True, color=CREAM_T)
txt(sl, 'Unidad 2 — Bases filosóficas del AC',
    ML, Cm(11.5), Cm(16), Cm(2), SANS, 17, color=SKY)
txt(sl, 'Lic. Jean Clemotte  |  PsicoEduca',
    ML, Cm(13.8), Cm(16), Cm(1.3), SANS, 15, italic=True, color=CREAM_T)

logo(sl, dark=True)


# ════════════════════════════════════════════════════════
# GUARDAR
# ════════════════════════════════════════════════════════
output = r'C:\Users\MI PC\psicoeduca\materiales\presentacion-unidad2.pptx'
prs.save(output)
print('Listo: ' + output)
print('18 slides | Unidad 2 | Rebranding | Lic. Jean Clemotte')
